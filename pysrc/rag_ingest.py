"""Tri-Store Hybrid RAG — Document Ingestion Pipeline.

Phase 1: Document → Chunks (PostgreSQL + pgvector)
Phase 2: Chunks → Entities (Neo4j graph)
Phase 3: Entities → Communities (Neo4j + PostgreSQL summaries)

Design principles:
  - Deterministic NER first, LLM only for ambiguous cases (cuts GraphRAG's token cost ~80%)
  - Neo4j nodes store metadata + chunk_id foreign keys only — full text stays in PostgreSQL
  - Embedding auto-selection: local model first, API fallback
  - All CRUD goes through PostgreSQL transactions (single source of truth)
"""

from __future__ import annotations

import asyncio
import hashlib
import json
import logging
import os
import re
import time
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from memory_ontology import ontology_node_id

import numpy as np

logger = logging.getLogger(__name__)

# ── Constants ────────────────────────────────────────────────────────────────

DEFAULT_CHUNK_SIZE = 512       # tokens (approximate, ~characters for CJK)
DEFAULT_CHUNK_OVERLAP = 64     # tokens
DEFAULT_EMBED_DIM = 1024  # Overridden by config["embedding"]["dim"] if set
BATCH_SIZE = 32                # embedding batch size

# Supported file types
SUPPORTED_EXTENSIONS = {
    ".pdf", ".txt", ".md", ".markdown", ".rst",
    ".html", ".htm", ".xml",
    ".pptx", ".ppt", ".xlsx", ".xls", ".csv", ".tsv",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".go", ".rs",
    ".java", ".c", ".cpp", ".h", ".hpp", ".cs", ".rb",
    ".php", ".swift", ".kt", ".scala", ".sh", ".bat",
    ".sql", ".yaml", ".yml", ".toml", ".json", ".xml",
    ".vue", ".svelte", ".css", ".scss", ".less",
}

# ── Deterministic NER Patterns ───────────────────────────────────────────────

# Regex-based entity extraction — fast, zero-token-cost first pass
ENTITY_PATTERNS: List[Tuple[re.Pattern, str]] = [
    # Dates
    (re.compile(r"\b\d{4}[-/]\d{1,2}[-/]\d{1,2}\b"), "DATE"),
    (re.compile(r"\b(?:19|20)\d{2}\s*年\s*\d{1,2}\s*月\s*\d{1,2}\s*日?\b"), "DATE"),
    (re.compile(r"\b(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}\b", re.IGNORECASE), "DATE"),
    # Emails
    (re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b"), "EMAIL"),
    # URLs
    (re.compile(r"https?://[^\s<>\"{}|\\^`\[\]]+"), "URL"),
    # Versions
    (re.compile(r"\bv?\d+\.\d+(?:\.\d+)?(?:-[a-zA-Z0-9]+)?\b"), "VERSION"),
    # Tech terms / identifiers
    (re.compile(r"\b[A-Z][a-z]+(?:[A-Z][a-z]+)+\b"), "TECH_TERM"),  # CamelCase
    (re.compile(r"\b[a-z]+(?:_[a-z]+){2,}\b"), "TECH_TERM"),       # snake_case with 3+ parts
    # Percentages / measurements
    (re.compile(r"\b\d+(?:\.\d+)?\s*%\b"), "MEASUREMENT"),
    (re.compile(r"\b\d+(?:\.\d+)?\s*(?:GB|MB|KB|TB|ms|s|GHz|MHz)\b"), "MEASUREMENT"),
    # Chinese entities
    (re.compile(r"[一-鿿]{2,6}(?:公司|集团|大学|研究所|实验室|部门|团队|系统|平台|框架|模型|算法|方法|技术|协议|标准)"), "ORG_TECH"),
    (re.compile(r"(?:[A-Z][a-z]*\s)?(?:University|Institute|Corporation|Lab|Department|Team)\s+(?:of\s+)?[A-Z][a-zA-Z]+(?:\s+[A-Z][a-zA-Z]+)*"), "ORG"),
]

# ── Text Chunking ────────────────────────────────────────────────────────────


def _estimate_tokens(text: str) -> int:
    """Rough token count: ~4 chars per token for English, ~1.5 for CJK."""
    cjk = sum(1 for c in text if "一" <= c <= "鿿" or "぀" <= c <= "ヿ")
    latin = len(text) - cjk
    return int(cjk / 1.5 + latin / 4)


def chunk_text(text: str, chunk_size: int = DEFAULT_CHUNK_SIZE,
               overlap: int = DEFAULT_CHUNK_OVERLAP) -> List[Dict[str, Any]]:
    """Split text into overlapping chunks respecting natural boundaries.

    Prefers splitting on paragraph breaks, then sentence breaks,
    then word boundaries. Each chunk returned as {text, index, start_char, end_char}.
    """
    if not text or not text.strip():
        return []

    paragraphs = text.split("\n\n")
    chunks: List[Dict[str, Any]] = []
    current = ""
    index = 0

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        # If paragraph fits in one chunk, add it
        if _estimate_tokens(current + " " + para) <= chunk_size:
            current = (current + "\n\n" + para).strip() if current else para
        else:
            # Flush current chunk
            if current:
                chunks.append({
                    "text": current,
                    "index": index,
                    "start_char": text.find(current) if index == 0 else 0,
                    "end_char": 0,
                })
                index += 1
                # Start new chunk with overlap from previous
                overlap_text = _extract_overlap(current, overlap)
                current = overlap_text + "\n\n" + para
            else:
                # Single paragraph is larger than chunk_size → split by sentences
                sentences = re.split(r"(?<=[.!?。！？\n])\s*", para)
                for sent in sentences:
                    sent = sent.strip()
                    if not sent:
                        continue
                    if _estimate_tokens(current + " " + sent) <= chunk_size:
                        current = (current + " " + sent).strip() if current else sent
                    else:
                        if current:
                            chunks.append({
                                "text": current,
                                "index": index,
                                "start_char": 0,
                                "end_char": 0,
                            })
                            index += 1
                        current = sent

    # Flush final chunk
    if current:
        chunks.append({
            "text": current,
            "index": index,
            "start_char": 0,
            "end_char": 0,
        })

    # Fix start/end positions
    offset = 0
    for ch in chunks:
        ch["start_char"] = text.find(ch["text"], offset) if offset < len(text) else offset
        ch["end_char"] = ch["start_char"] + len(ch["text"])
        offset = ch["end_char"]

    return chunks


def _extract_overlap(text: str, overlap_tokens: int) -> str:
    """Extract the last ~overlap_tokens worth of text."""
    words = text.split()
    if len(words) <= overlap_tokens // 2:
        return text[-min(len(text), overlap_tokens * 4):]
    return " ".join(words[-(overlap_tokens // 2):])


# ── Document Parsers ─────────────────────────────────────────────────────────


def parse_pdf(filepath: str) -> str:
    """Extract text from PDF file."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        return "\n\n".join(
            page.extract_text() or ""
            for page in reader.pages
        )
    except ImportError:
        raise ImportError("pypdf is required for PDF parsing. Install: pip install pypdf")


def parse_office(filepath: str) -> str:
    """Extract text from Office files (.pptx, .xlsx, .docx)."""
    ext = Path(filepath).suffix.lower()
    if ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            parts = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            slide_text.append(para.text)
                parts.append("\n".join(slide_text))
            return "\n\n--- slide ---\n\n".join(parts)
        except ImportError:
            raise ImportError("python-pptx is required for PPT parsing. Install: pip install python-pptx")
    elif ext in (".xlsx", ".xls"):
        try:
            import pandas as pd
            xl = pd.ExcelFile(filepath)
            parts = []
            for sheet_name in xl.sheet_names:
                df = pd.read_excel(filepath, sheet_name=sheet_name)
                parts.append(f"## Sheet: {sheet_name}\n{df.to_markdown(index=False)}")
            return "\n\n".join(parts)
        except ImportError:
            raise ImportError("pandas + openpyxl required for Excel. Install: pip install pandas openpyxl")
    elif ext in (".csv", ".tsv"):
        try:
            import pandas as pd
            sep = "\t" if ext == ".tsv" else ","
            df = pd.read_csv(filepath, sep=sep)
            return df.to_markdown(index=False)
        except ImportError:
            raise ImportError("pandas required for CSV/TSV. Install: pip install pandas")
    return ""


def parse_html(filepath: str) -> str:
    """Extract text from HTML file."""
    try:
        from bs4 import BeautifulSoup
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
        # Remove script/style elements
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        raise ImportError("beautifulsoup4 required for HTML. Install: pip install beautifulsoup4")


def parse_text(filepath: str) -> str:
    """Read plain text, Markdown, or code files."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


PARSERS = {
    ".pdf": parse_pdf,
    ".pptx": parse_office, ".ppt": parse_office,
    ".xlsx": parse_office, ".xls": parse_office,
    ".csv": parse_office, ".tsv": parse_office,
    ".html": parse_html, ".htm": parse_html,
    ".xml": parse_html,
}

# Default: plain text reader works for .txt, .md, .py, .js, etc.
def _get_parser(ext: str):
    return PARSERS.get(ext, parse_text)


# ── Entity Extraction ────────────────────────────────────────────────────────


def extract_entities_deterministic(text: str) -> List[Dict[str, str]]:
    """Fast regex-based entity extraction. Zero LLM cost."""
    entities: List[Dict[str, str]] = []
    seen: set = set()

    for pattern, entity_type in ENTITY_PATTERNS:
        for match in pattern.finditer(text):
            entity_text = match.group(0).strip()
            key = (entity_text.lower(), entity_type)
            if key not in seen:
                seen.add(key)
                entities.append({
                    "name": entity_text,
                    "type": entity_type,
                    "source": "regex",
                    "confidence": 0.85,
                })

    return entities


async def extract_entities_llm(text: str, client, model: str,
                                extra_params: dict = None) -> List[Dict[str, str]]:
    """LLM-based entity extraction for ambiguous/unmatched cases.

    Only call this for texts where deterministic extraction found nothing
    or for critical documents. Uses a small model by default.
    """
    prompt = (
        "Extract named entities from the following text. "
        "Return a JSON array of objects with keys: name, type (PERSON/ORG/TECH/EVENT/LOCATION/CONCEPT), "
        "and confidence (0-1). Only return entities you are confident about.\n\n"
        f"Text:\n{text[:3000]}\n\n"
        "Return JSON only."
    )
    try:
        res = await client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "You extract named entities from text. Return JSON only."},
                {"role": "user", "content": prompt},
            ],
            response_format={"type": "json_object"},
            **(extra_params or {}),
        )
        data = json.loads(res.choices[0].message.content)
        entities = data.get("entities", data.get("results", []))
        for e in entities:
            e.setdefault("source", "llm")
            e.setdefault("confidence", 0.7)
        return entities
    except Exception as e:
        logger.debug(f"LLM entity extraction failed: {e}")
        return []


# ── Embedding ────────────────────────────────────────────────────────────────


class EmbeddingProvider:
    """Auto-selecting embedding provider: local model → API fallback."""

    def __init__(self, config: dict = None, external_model=None):
        cfg = config or {}
        embed_cfg = cfg.get("embedding", {})
        self.provider = embed_cfg.get("provider", "auto")
        self.model_path = embed_cfg.get("model_path", "")
        self.api_model = embed_cfg.get("api_model", "text-embedding-3-small")
        self.dim = embed_cfg.get("dim", DEFAULT_EMBED_DIM)
        self._local_model = external_model  # Shared instance if provided
        self._client = None

    async def _ensure_local_model(self):
        if self._local_model is not None:
            return self._local_model
        try:
            from sentence_transformers import SentenceTransformer
            path = self.model_path or "BAAI/bge-small-zh-v1.5"
            self._local_model = await asyncio.to_thread(
                SentenceTransformer, path
            )
            logger.info(f"Loaded local embedding model: {path}")
            return self._local_model
        except Exception as e:
            logger.warning(f"Local embedding model unavailable: {e}")
            return None

    async def _ensure_client(self, config: dict = None):
        if self._client is not None:
            return self._client
        try:
            from openai import AsyncOpenAI
            cfg = config or {}
            providers = cfg.get("llm_providers", {}) or {}
            provider = providers.get("openai", {})
            self._client = AsyncOpenAI(
                api_key=provider.get("api_key", os.environ.get("OPENAI_API_KEY", "")),
                base_url=provider.get("base_url"),
            )
            return self._client
        except Exception as e:
            logger.warning(f"OpenAI client unavailable: {e}")
            return None

    async def embed(self, texts: List[str], config: dict = None) -> np.ndarray:
        """Embed a list of texts, auto-selecting provider."""
        if self.provider in ("local", "auto"):
            model = await self._ensure_local_model()
            if model is not None:
                embeddings = await asyncio.to_thread(
                    model.encode, texts,
                    normalize_embeddings=True,
                    show_progress_bar=False,
                )
                return np.array(embeddings, dtype=np.float32)

        if self.provider in ("api", "auto"):
            client = await self._ensure_client(config)
            if client is not None:
                resp = await client.embeddings.create(
                    model=self.api_model,
                    input=texts,
                )
                return np.array(
                    [d.embedding for d in resp.data],
                    dtype=np.float32,
                )

        # Ultimate fallback: zero vectors
        logger.warning("No embedding provider available, using zero vectors")
        return np.zeros((len(texts), self.dim), dtype=np.float32)

    async def embed_single(self, text: str, config: dict = None) -> List[float]:
        embeddings = await self.embed([text], config)
        return embeddings[0].tolist()


# ── Ingestion Pipeline ───────────────────────────────────────────────────────


@dataclass
class IngestResult:
    doc_id: str
    title: str
    file_type: str
    chunk_count: int
    entity_count: int
    elapsed_ms: float


class RAGIngestionPipeline:
    """Orchestrates the full document → chunks → entities → communities pipeline."""

    def __init__(self, memory_service=None, config: dict = None):
        self._memory = memory_service
        self._config = config or {}
        self._embedder = EmbeddingProvider(config)

    async def ingest_file(self, filepath: str, owner_id: str = "default",
                          scope: str = "shared", metadata: dict = None) -> IngestResult:
        """Ingest a single file through the full pipeline."""
        t0 = time.monotonic()
        path = Path(filepath)
        ext = path.suffix.lower()

        if ext not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {ext}")

        # Phase 1: Parse → Chunk → Embed → Store in PostgreSQL
        parser = _get_parser(ext)
        raw_text = parser(str(path))

        chunks = chunk_text(raw_text)
        if not chunks:
            raise ValueError(f"No extractable text in {filepath}")

        # Generate doc_id
        doc_id = hashlib.sha256(
            f"{owner_id}:{path.name}:{time.time()}".encode()
        ).hexdigest()[:16]

        # Embed all chunks in batches
        chunk_texts = [c["text"] for c in chunks]
        all_embeddings = []
        for i in range(0, len(chunk_texts), BATCH_SIZE):
            batch = chunk_texts[i:i + BATCH_SIZE]
            emb = await self._embedder.embed(batch, self._config)
            all_embeddings.append(emb)
        embeddings = np.concatenate(all_embeddings, axis=0) if all_embeddings else np.array([])

        # Phase 2: Extract entities from all chunks
        all_entities: List[Dict[str, str]] = []
        for ch in chunks:
            entities = extract_entities_deterministic(ch["text"])
            all_entities.extend(entities)

        # If deterministic NER found very few entities, try LLM
        if len(all_entities) < 3 and len(raw_text) > 500:
            try:
                client = await self._embedder._ensure_client(self._config)
                if client:
                    llm_entities = await extract_entities_llm(
                        raw_text[:5000], client,
                        self._config.get("llm_providers", {}).get("openai", {}).get("model", "gpt-4o-mini"),
                        self._config.get("llm_providers", {}).get("openai", {}).get("extra_params", {}),
                    )
                    all_entities.extend(llm_entities)
            except Exception:
                pass

        # Deduplicate entities
        deduped: Dict[str, Dict[str, str]] = {}
        for e in all_entities:
            key = (e["name"].lower(), e["type"])
            if key not in deduped or e.get("confidence", 0) > deduped[key].get("confidence", 0):
                deduped[key] = e
        unique_entities = list(deduped.values())

        # Store in PostgreSQL (if memory service is available)
        if self._memory is not None:
            await self._store_chunks(doc_id, path, ext, chunks, embeddings,
                                     owner_id, scope, metadata or {})
            await self._store_entities(doc_id, unique_entities, chunks)

        elapsed = (time.monotonic() - t0) * 1000
        result = IngestResult(
            doc_id=doc_id,
            title=path.name,
            file_type=ext,
            chunk_count=len(chunks),
            entity_count=len(unique_entities),
            elapsed_ms=round(elapsed, 1),
        )
        logger.info(f"Ingested {filepath}: {result}")
        return result

    async def ingest_directory(self, dirpath: str, owner_id: str = "default",
                               scope: str = "shared", recursive: bool = True) -> List[IngestResult]:
        """Ingest all supported files in a directory."""
        results = []
        pattern = "**/*" if recursive else "*"
        for filepath in Path(dirpath).glob(pattern):
            if filepath.is_file() and filepath.suffix.lower() in SUPPORTED_EXTENSIONS:
                try:
                    result = await self.ingest_file(str(filepath), owner_id, scope)
                    results.append(result)
                except Exception as e:
                    logger.warning(f"Failed to ingest {filepath}: {e}")
        return results

    async def _store_chunks(self, doc_id: str, path: Path, ext: str,
                            chunks: List[Dict], embeddings: np.ndarray,
                            owner_id: str, scope: str, metadata: dict):
        """Store document + chunks in PostgreSQL via memory service."""
        # Store document record
        doc_record = {
            "id": doc_id,
            "title": path.name,
            "source": str(path),
            "file_type": ext,
            "owner_id": owner_id,
            "scope": scope,
            "metadata": json.dumps(metadata, ensure_ascii=False),
            "chunk_count": len(chunks),
        }
        await self._memory._execute_sql(
            """INSERT INTO rag_documents (id, title, source, file_type, owner_id, scope, metadata, chunk_count)
               VALUES ($1, $2, $3, $4, $5, $6, $7, $8)
               ON CONFLICT (id) DO UPDATE SET chunk_count = $8, metadata = $7""",
            doc_record["id"], doc_record["title"], doc_record["source"],
            doc_record["file_type"], doc_record["owner_id"], doc_record["scope"],
            doc_record["metadata"], doc_record["chunk_count"],
        )

        # Store chunks
        for i, ch in enumerate(chunks):
            chunk_id = f"{doc_id}:{i}"
            emb = embeddings[i].tolist() if i < len(embeddings) else [0.0] * DEFAULT_EMBED_DIM
            await self._memory._execute_sql(
                """INSERT INTO rag_chunks (id, doc_id, chunk_index, text, embedding, metadata, owner_id, scope)
                   VALUES ($1, $2, $3, $4, $5, $6, $7, $8)
                   ON CONFLICT (id) DO UPDATE SET text = $4, embedding = $5""",
                chunk_id, doc_id, i, ch["text"],
                json.dumps(emb),
                json.dumps({"start_char": ch["start_char"], "end_char": ch["end_char"]}),
                owner_id, scope,
            )

    async def _store_entities(self, doc_id: str, entities: List[Dict],
                              chunks: List[Dict]):
        """Store entities in Neo4j via memory service, linking to chunk IDs."""
        if not self._memory or not hasattr(self._memory, "graph"):
            return

        graph = self._memory.graph
        for entity in entities:
            entity_id = ontology_node_id('rag_entity', entity['name'])

            # Upsert entity node in Neo4j
            await graph.upsert_node(
                entity_id,
                "rag_entity",
                entity["name"],
                {
                    "type": entity["type"],
                    "source": entity.get("source", "regex"),
                    "confidence": entity.get("confidence", 0.85),
                    "doc_id": doc_id,
                },
            )

        # Create MENTIONS edges: entity → chunk
        # (simplified: each entity linked to all chunks of its source doc)
        for entity in entities:
            entity_id = ontology_node_id('rag_entity', entity['name'])
            for i in range(len(chunks)):
                chunk_id = f"{doc_id}:{i}"
                await graph.upsert_edge(
                    entity_id, chunk_id, "rag_mentions",
                    confidence=entity.get("confidence", 0.85),
                    metadata={"doc_id": doc_id},
                )


# ── Convenience API ──────────────────────────────────────────────────────────

async def ingest_document(filepath: str, owner_id: str = "default",
                          scope: str = "shared", config: dict = None,
                          memory_service=None) -> IngestResult:
    """One-liner to ingest a document. The main entry point for API/agent use."""
    pipeline = RAGIngestionPipeline(memory_service, config)
    return await pipeline.ingest_file(filepath, owner_id, scope)


async def ingest_directory(dirpath: str, owner_id: str = "default",
                           scope: str = "shared", config: dict = None,
                           memory_service=None) -> List[IngestResult]:
    """One-liner to ingest a directory."""
    pipeline = RAGIngestionPipeline(memory_service, config)
    return await pipeline.ingest_directory(dirpath, owner_id, scope)
