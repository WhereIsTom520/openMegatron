"""
PPT Master Enhanced Layouts - 增强版布局系统

新增 5 种专业布局:
1. code_block    - 代码块展示（深色主题，等宽字体）
2. flowchart     - 流程/步骤图（编号圆圈+箭头连接）
3. three_column  - 三栏对比布局
4. grid_2x2      - 2x2 卡片网格
5. comparison_plus - 增强对比（✓/✗ 标记，特性矩阵）
"""

from __future__ import annotations
from typing import Any
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
import re


def _hex_to_rgb(value: str) -> RGBColor:
    text = str(value or "").strip().lstrip("#")
    if len(text) == 3:
        text = "".join(ch * 2 for ch in text)
    if not re.fullmatch(r"[0-9a-fA-F]{6}", text):
        text = "2563eb"
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def _text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, (list, tuple)):
        return "\n".join(str(item) for item in value)
    return str(value)


def _as_bullets(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, str):
        return [re.sub(r"^\s*[-*+]\s+", "", line).strip()
                for line in value.splitlines() if line.strip()]
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()]
    return [str(value).strip()] if str(value).strip() else []


def _make_font(size: int, color: str, bold: bool = False, italic: bool = False,
               name: str = None) -> dict:
    return {"size": Pt(size), "color": color, "bold": bold, "italic": italic, "name": name}


def _apply_font(run, font_spec: dict, theme: dict) -> None:
    run.font.size = font_spec.get("size", Pt(16))
    run.font.bold = font_spec.get("bold", False)
    if font_spec.get("italic"):
        run.font.italic = True
    run.font.color.rgb = _hex_to_rgb(font_spec.get("color", theme["dark"]))
    name = font_spec.get("name") or theme.get("font_body")
    if name:
        run.font.name = name


def _add_textbox(slide, left, top, width, height, text: str,
                 font_spec: dict, theme: dict, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, word_wrap: bool = True,
                 line_spacing: float = 1.15) -> Any:
    box = slide.shapes.add_textbox(
        Inches(left) if isinstance(left, (int, float)) else left,
        Inches(top) if isinstance(top, (int, float)) else top,
        Inches(width) if isinstance(width, (int, float)) else width,
        Inches(height) if isinstance(height, (int, float)) else height,
    )
    frame = box.text_frame
    frame.word_wrap = word_wrap
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.alignment = align
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    _apply_font(run, font_spec, theme)
    return box


def _add_bullets(slide, bullets: list[str], left, top, width, height,
                 font_spec: dict, theme: dict, line_spacing: float = 1.3) -> Any:
    box = slide.shapes.add_textbox(
        Inches(left) if isinstance(left, (int, float)) else left,
        Inches(top) if isinstance(top, (int, float)) else top,
        Inches(width) if isinstance(width, (int, float)) else width,
        Inches(height) if isinstance(height, (int, float)) else height,
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.08)
    for i, bullet in enumerate(bullets):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.text = bullet
        para.level = 0
        para.space_after = Pt(6)
        run = para.runs[0] if para.runs else para.add_run()
        _apply_font(run, font_spec, theme)
    return box


def _add_shape_bg(slide, left, top, width, height, color: str) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left) if isinstance(left, (int, float)) else left,
        Inches(top) if isinstance(top, (int, float)) else top,
        Inches(width) if isinstance(width, (int, float)) else width,
        Inches(height) if isinstance(height, (int, float)) else height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(color)
    shape.line.fill.background()
    return shape


def _add_connector_line(slide, x1, y1, x2, y2, color: str, width: float = 2):
    """Draw a thin rectangle as a connector line (MSO_SHAPE.LINE not available in all versions)."""
    if abs(x1 - x2) < 0.01:
        # Vertical line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x1 - 0.01), Inches(min(y1, y2)),
            Inches(0.02), Inches(abs(y1 - y2)),
        )
    else:
        # Horizontal line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(min(x1, x2)), Inches(y1 - 0.01),
            Inches(abs(x1 - x2)), Inches(0.02),
        )
    line.fill.solid()
    line.fill.fore_color.rgb = _hex_to_rgb(color)
    line.line.fill.background()
    return line


def _add_accent_bar(slide, theme: dict):
    """统一装饰条"""
    _add_shape_bg(slide, 0.6, 0.35, 0.06, 0.55, theme["accent"])


# ═══════════════════════════════════════════════════════════════
#  1. CODE BLOCK LAYOUT
# ═══════════════════════════════════════════════════════════════

def build_code_block(prs: Presentation, item: dict, index: int, theme: dict) -> None:
    """代码块布局 - 深色主题编辑器风格

    item fields:
        title: slide title
        code: source code string
        language: language identifier (python, javascript, etc.)
        caption: optional caption below code
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide, theme)

    title = str(item.get("title") or f"Slide {index}")
    _add_textbox(slide, 0.78, 0.32, 11.5, 0.7, title,
                 _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    code = str(item.get("code") or item.get("body") or "")
    language = str(item.get("language") or "").lower()

    # 代码背景
    code_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.2),
    )
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    code_bg.line.color.rgb = RGBColor(0x44, 0x44, 0x55)
    code_bg.line.width = Pt(1)

    # 语言标签
    if language:
        lang_badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.0), Inches(1.42), Inches(1.2), Inches(0.32),
        )
        lang_badge.fill.solid()
        lang_badge.fill.fore_color.rgb = _hex_to_rgb(theme["accent"])
        lang_badge.line.fill.background()
        _add_textbox(slide, 1.0, 1.40, 1.2, 0.32, language.upper(),
                     _make_font(9, "#ffffff", bold=True), theme,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 代码文本
    display_code = code[:1500]
    if len(code) > 1500:
        display_code += "\n# ... (truncated)"

    text_frame = slide.shapes.add_textbox(
        Inches(1.2), Inches(1.95), Inches(10.8), Inches(4.2),
    ).text_frame
    text_frame.word_wrap = True

    lines = display_code.split('\n')
    keywords_python = {"def", "class", "import", "from", "return", "if", "else", "elif",
                       "for", "while", "try", "except", "with", "as", "in", "not", "and", "or",
                       "True", "False", "None", "yield", "raise", "pass", "break", "continue"}
    keywords_js = {"function", "const", "let", "var", "return", "if", "else", "for", "while",
                   "try", "catch", "class", "new", "this", "async", "await", "export", "import"}

    for i, line in enumerate(lines):
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        para.text = line
        para.font.name = "Consolas"
        para.font.size = Pt(11)

        # 基础着色
        stripped = line.strip()
        if stripped.startswith("#") or stripped.startswith("//"):
            color = RGBColor(0x6A, 0x99, 0x5A)  # 绿色注释
        elif stripped.startswith('"""') or stripped.startswith("'''"):
            color = RGBColor(0xCE, 0x91, 0x78)  # 橙色文档字符串
        else:
            first_word = stripped.split()[0] if stripped else ""
            if first_word.rstrip(":") in keywords_python or first_word.rstrip(":") in keywords_js:
                color = RGBColor(0x56, 0x9C, 0xD6)  # 蓝色关键字
            elif first_word in ("def", "class", "function"):
                color = RGBColor(0xAE, 0xF3, 0xC4)  # 绿色函数/类
            else:
                color = RGBColor(0xD4, 0xD4, 0xD4)  # 灰色默认

        para.font.color.rgb = color
        para.space_after = Pt(2)

    # 标题
    caption = str(item.get("caption") or "")
    if caption:
        _add_textbox(slide, 0.8, 6.6, 11.5, 0.4, caption,
                     _make_font(11, theme["muted"]), theme, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
#  2. FLOWCHART LAYOUT
# ═══════════════════════════════════════════════════════════════

def build_flowchart(prs: Presentation, item: dict, index: int, theme: dict) -> None:
    """流程图/步骤图布局

    item fields:
        title: slide title
        steps: list of {title, description} dicts
        direction: "vertical" (default) or "horizontal"
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide, theme)

    title = str(item.get("title") or f"Slide {index}")
    _add_textbox(slide, 0.78, 0.32, 11.5, 0.7, title,
                 _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    steps = item.get("steps") or item.get("process") or item.get("bullets") or []
    if not isinstance(steps, list) or len(steps) == 0:
        _add_textbox(slide, 0.8, 3.0, 11.5, 1.0, "No steps provided",
                     _make_font(16, theme["muted"]), theme, align=PP_ALIGN.CENTER)
        return

    direction = str(item.get("direction") or "vertical").lower()

    if direction == "horizontal":
        _build_flowchart_horizontal(slide, steps, theme)
    else:
        _build_flowchart_vertical(slide, steps, theme)


def _build_flowchart_vertical(slide, steps: list, theme: dict) -> None:
    """垂直流程图 - 适用于 <=6 个步骤"""
    n = min(len(steps), 6)
    box_w = 10.0
    box_h = 0.85
    start_y = 1.5

    for i in range(n):
        step = steps[i]
        step_title = ""
        step_desc = ""
        if isinstance(step, dict):
            step_title = str(step.get("title") or step.get("label") or f"Step {i + 1}")
            step_desc = str(step.get("description") or step.get("text") or step.get("desc") or "")
        else:
            step_title = str(step)

        y = start_y + i * (box_h + 0.45)

        # 步骤编号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.1), Inches(0.65), Inches(0.65),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = _hex_to_rgb(theme["accent"])
        circle.line.fill.background()
        _add_textbox(slide, 0.8, y + 0.1, 0.65, 0.65, str(i + 1),
                     _make_font(18, "#ffffff", bold=True, name=theme.get("font_title")),
                     theme, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 步骤框
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(y),
            Inches(box_w), Inches(box_h),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex_to_rgb(theme["light"])
        rect.line.color.rgb = _hex_to_rgb(theme["accent"])
        rect.line.width = Pt(1.5)

        # 标题
        _add_textbox(slide, 2.0, y + 0.05, box_w - 0.4, 0.32, step_title,
                     _make_font(14, theme["dark"], bold=True, name=theme.get("font_body")),
                     theme)

        # 描述
        if step_desc:
            _add_textbox(slide, 2.0, y + 0.42, box_w - 0.4, 0.35, step_desc,
                         _make_font(10, theme["muted"], name=theme.get("font_body")),
                         theme)

        # 连接箭头
        if i < n - 1:
            arrow_y = y + box_h
            _add_connector_line(slide, 1.125, arrow_y, 1.125, arrow_y + 0.45, theme["accent"])
            # 箭头尖
            tip = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(1.0), Inches(arrow_y + 0.3),
                Inches(0.25), Inches(0.25),
            )
            tip.fill.solid()
            tip.fill.fore_color.rgb = _hex_to_rgb(theme["accent"])
            tip.line.fill.background()
            tip.rotation = 180


def _build_flowchart_horizontal(slide, steps: list, theme: dict) -> None:
    """水平流程图 - 适用于 <=5 个步骤"""
    n = min(len(steps), 5)
    box_w = 2.0
    box_h = 2.2
    arrow_w = 0.6
    total_w = n * box_w + (n - 1) * arrow_w
    start_x = (13.333 - total_w) / 2
    start_y = 2.5

    for i in range(n):
        step = steps[i]
        step_title = ""
        step_desc = ""
        if isinstance(step, dict):
            step_title = str(step.get("title") or step.get("label") or f"{i + 1}")
            step_desc = str(step.get("description") or step.get("text") or "")
        else:
            step_title = str(step)

        x = start_x + i * (box_w + arrow_w)

        # 步骤框
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(start_y),
            Inches(box_w), Inches(box_h),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex_to_rgb(theme["light"])
        rect.line.color.rgb = _hex_to_rgb(theme["accent"])
        rect.line.width = Pt(2)

        # 编号
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + box_w / 2 - 0.3), Inches(start_y + 0.15),
            Inches(0.6), Inches(0.6),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = _hex_to_rgb(theme["accent"])
        circle.line.fill.background()
        _add_textbox(slide, x + box_w / 2 - 0.3, start_y + 0.15, 0.6, 0.6, str(i + 1),
                     _make_font(16, "#ffffff", bold=True), theme,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 标题
        _add_textbox(slide, x + 0.1, start_y + 0.95, box_w - 0.2, 0.4, step_title,
                     _make_font(11, theme["dark"], bold=True, name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER)

        # 描述
        if step_desc:
            _add_textbox(slide, x + 0.1, start_y + 1.4, box_w - 0.2, 0.6, step_desc,
                         _make_font(9, theme["muted"], name=theme.get("font_body")),
                         theme, align=PP_ALIGN.CENTER)

        # 箭头
        if i < n - 1:
            arrow_x = x + box_w
            arrow_center_y = start_y + box_h / 2
            _add_connector_line(slide, arrow_x, arrow_center_y, arrow_x + arrow_w - 0.1, arrow_center_y, theme["accent"])
            # 箭头尖
            tip = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE, Inches(arrow_x + arrow_w - 0.3),
                Inches(arrow_center_y - 0.12),
                Inches(0.2), Inches(0.24),
            )
            tip.fill.solid()
            tip.fill.fore_color.rgb = _hex_to_rgb(theme["accent"])
            tip.line.fill.background()
            tip.rotation = 90


# ═══════════════════════════════════════════════════════════════
#  3. THREE COLUMN LAYOUT
# ═══════════════════════════════════════════════════════════════

def build_three_column(prs: Presentation, item: dict, index: int, theme: dict) -> None:
    """三栏布局 - 功能对比、选项对比

    item fields:
        title: slide title
        column1/column2/column3: lists of bullets
        title1/title2/title3: column headers
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide, theme)

    title = str(item.get("title") or f"Slide {index}")
    _add_textbox(slide, 0.78, 0.32, 11.5, 0.7, title,
                 _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    cols = [
        item.get("column1") or item.get("left") or [],
        item.get("column2") or item.get("center") or item.get("middle") or [],
        item.get("column3") or item.get("right") or [],
    ]

    col_titles = [
        str(item.get("title1") or item.get("title_left") or "Option A"),
        str(item.get("title2") or item.get("title_center") or "Option B"),
        str(item.get("title3") or item.get("title_right") or "Option C"),
    ]

    col_colors = [
        theme["accent"],
        theme.get("secondary", "#14b8a6"),
        theme.get("success", "#10b981"),
    ]
    col_width = 3.6

    for col_idx, (col_data, col_title, color) in enumerate(zip(cols, col_titles, col_colors)):
        x = 0.8 + col_idx * (col_width + 0.25)

        # 列标题
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3),
            Inches(col_width), Inches(0.55),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = _hex_to_rgb(color)
        header.line.fill.background()
        _add_textbox(slide, x, 1.28, col_width, 0.55, col_title,
                     _make_font(14, "#ffffff", bold=True, name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 列内容框
        body = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9),
            Inches(col_width), Inches(4.8),
        )
        body.fill.solid()
        body.fill.fore_color.rgb = _hex_to_rgb(theme["light"])
        body.line.color.rgb = _hex_to_rgb(color)
        body.line.width = Pt(1)

        # 列内容
        bullets = _as_bullets(col_data)
        if bullets:
            text_frame = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(2.1), Inches(col_width - 0.4), Inches(4.4),
            ).text_frame
            text_frame.word_wrap = True
            for i, bullet in enumerate(bullets[:10]):
                para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                para.text = f"  {bullet}"
                para.font.size = Pt(11)
                para.font.color.rgb = _hex_to_rgb(theme["dark"])
                para.space_after = Pt(8)


# ═══════════════════════════════════════════════════════════════
#  4. GRID 2x2 LAYOUT
# ═══════════════════════════════════════════════════════════════

def build_grid_2x2(prs: Presentation, item: dict, index: int, theme: dict,
                   image_placeholders: bool = True) -> None:
    """2x2 卡片网格布局

    item fields:
        title: slide title
        cards: list of {title, description, icon} dicts (max 4)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide, theme)

    title = str(item.get("title") or f"Slide {index}")
    _add_textbox(slide, 0.78, 0.32, 11.5, 0.7, title,
                 _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    cards = item.get("cards") or item.get("items") or item.get("bullets") or []
    if not isinstance(cards, list):
        cards = []

    card_w = 5.5
    card_h = 2.5
    icon_colors = [theme["accent"], theme.get("secondary", "#14b8a6"),
                   theme.get("success", "#10b981"), theme.get("warning", "#f59e0b")]

    for card_idx in range(min(4, len(cards))):
        row = card_idx // 2
        col = card_idx % 2
        x = 0.8 + col * (card_w + 0.4)
        y = 1.3 + row * (card_h + 0.35)

        card_data = cards[card_idx]
        card_title = ""
        card_desc = ""
        card_icon = None

        if isinstance(card_data, dict):
            card_title = str(card_data.get("title") or card_data.get("label") or f"Item {card_idx + 1}")
            card_desc = str(card_data.get("description") or card_data.get("text") or card_data.get("desc") or "")
            card_icon = card_data.get("icon") or card_data.get("emoji")
        else:
            card_title = str(card_data)

        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(card_w), Inches(card_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _hex_to_rgb(theme["light"])
        card.line.color.rgb = _hex_to_rgb(icon_colors[card_idx])
        card.line.width = Pt(1.5)

        # 图标圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.25), Inches(y + 0.25),
            Inches(0.55), Inches(0.55),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = _hex_to_rgb(icon_colors[card_idx])
        circle.line.fill.background()

        icon_text = str(card_icon or card_idx + 1)
        _add_textbox(slide, x + 0.25, y + 0.22, 0.55, 0.55, icon_text,
                     _make_font(14, "#ffffff", bold=True, name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 卡片标题
        _add_textbox(slide, x + 1.0, y + 0.28, card_w - 1.3, 0.4, card_title,
                     _make_font(14, theme["dark"], bold=True, name=theme.get("font_body")),
                     theme)

        # 卡片描述
        if card_desc:
            _add_textbox(slide, x + 0.25, y + 0.9, card_w - 0.5, 1.4, card_desc,
                         _make_font(10, theme["muted"], name=theme.get("font_body")),
                         theme)


# ═══════════════════════════════════════════════════════════════
#  5. COMPARISON PLUS LAYOUT
# ═══════════════════════════════════════════════════════════════

def build_comparison_plus(prs: Presentation, item: dict, index: int, theme: dict) -> None:
    """增强对比布局 - 带 ✓/✗ 标记的特性矩阵

    item fields:
        title: slide title
        left_title/right_title: column headers
        features: list of {name, left, right} dicts
            left/right values: True/False/string
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide, theme)

    title = str(item.get("title") or f"Slide {index}")
    _add_textbox(slide, 0.78, 0.32, 11.5, 0.7, title,
                 _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    left_title = str(item.get("left_title") or item.get("title_a") or "Option A")
    right_title = str(item.get("right_title") or item.get("title_b") or "Option B")
    color_a = theme["accent"]
    color_b = theme.get("secondary", "#14b8a6")

    # 列标题
    for col_idx, (col_title, color) in enumerate([(left_title, color_a), (right_title, color_b)]):
        x = 0.8 + col_idx * 6.2
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.2),
            Inches(5.8), Inches(0.55),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = _hex_to_rgb(color)
        header.line.fill.background()
        _add_textbox(slide, x, 1.18, 5.8, 0.55, col_title,
                     _make_font(15, "#ffffff", bold=True, name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 特性行
    features = item.get("features") or []
    if not features:
        left_items = _as_bullets(item.get("left_items") or item.get("items_a") or [])
        right_items = _as_bullets(item.get("right_items") or item.get("items_b") or [])
        if left_items:
            _add_bullets(slide, left_items[:8], 0.8, 2.0, 5.5, 4.8,
                         _make_font(13, theme["dark"], name=theme.get("font_body")), theme)
        if right_items:
            _add_bullets(slide, right_items[:8], 7.0, 2.0, 5.5, 4.8,
                         _make_font(13, theme["dark"], name=theme.get("font_body")), theme)
        return

    for i, feature in enumerate(features[:7]):
        y = 1.95 + i * 0.65
        feature_name = ""
        left_val = None
        right_val = None

        if isinstance(feature, dict):
            feature_name = str(feature.get("name") or str(feature))
            left_val = feature.get("left") or feature.get("a")
            right_val = feature.get("right") or feature.get("b")
        else:
            feature_name = str(feature)

        # 分隔线
        if i > 0:
            _add_connector_line(slide, 0.8, y - 0.05, 12.1, y - 0.05, "#E5E7EB", 0.5)

        # 特性名
        _add_textbox(slide, 0.8, y + 0.08, 2.8, 0.45, feature_name,
                     _make_font(12, theme["dark"], bold=True, name=theme.get("font_body")),
                     theme)

        # 左值
        _render_comparison_value(slide, left_val, 3.8, y, theme)

        # 分隔竖线
        _add_connector_line(slide, 6.65, y - 0.05, 6.65, y + 0.6, "#D1D5DB", 1.0)

        # 右值
        _render_comparison_value(slide, right_val, 7.0, y, theme)


def _render_comparison_value(slide, value, x, y, theme: dict) -> None:
    """渲染单个对比值（✓/✗ 或文本）"""
    positive_values = {True, "yes", "Yes", "true", "True", "✓", "✔", "√", "check"}
    negative_values = {False, "no", "No", "false", "False", "✗", "✘", "×", "x"}

    if value in positive_values:
        _add_textbox(slide, x, y - 0.05, 2.5, 0.5, "✓",
                     _make_font(22, "#10b981", bold=True), theme,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    elif value in negative_values:
        _add_textbox(slide, x, y - 0.05, 2.5, 0.5, "✗",
                     _make_font(22, "#ef4444", bold=True), theme,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    elif value is not None:
        _add_textbox(slide, x, y + 0.08, 2.5, 0.45, str(value),
                     _make_font(12, theme["muted"], name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER)
