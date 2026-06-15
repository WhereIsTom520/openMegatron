from __future__ import annotations

import json
import os
import re
import subprocess
import sys
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn

# ═══════════════════════════════════════════════════════════════
#  THEMES — expanded with more color palettes and typography
# ═══════════════════════════════════════════════════════════════

THEMES = {
    "professional": {
        "accent": "#2563eb", "dark": "#111827", "muted": "#64748b",
        "light": "#eff6ff", "background": "#ffffff",
        "secondary": "#0891b2", "success": "#059669", "warning": "#d97706",
        "font_title": "Calibri", "font_body": "Calibri",
    },
    "academic": {
        "accent": "#0f766e", "dark": "#172554", "muted": "#475569",
        "light": "#ecfeff", "background": "#ffffff",
        "secondary": "#6366f1", "success": "#15803d", "warning": "#b45309",
        "font_title": "Georgia", "font_body": "Georgia",
    },
    "creative": {
        "accent": "#db2777", "dark": "#1f2937", "muted": "#6b7280",
        "light": "#fdf2f8", "background": "#fff7ed",
        "secondary": "#8b5cf6", "success": "#16a34a", "warning": "#ea580c",
        "font_title": "Montserrat", "font_body": "Segoe UI",
    },
    "minimal": {
        "accent": "#111827", "dark": "#111827", "muted": "#71717a",
        "light": "#f4f4f5", "background": "#ffffff",
        "secondary": "#52525b", "success": "#4d7c0f", "warning": "#a16207",
        "font_title": "Helvetica", "font_body": "Helvetica",
    },
    "dark": {
        "accent": "#60a5fa", "dark": "#f8fafc", "muted": "#94a3b8",
        "light": "#1e293b", "background": "#0f172a",
        "secondary": "#818cf8", "success": "#34d399", "warning": "#fbbf24",
        "font_title": "Calibri", "font_body": "Calibri",
    },
    "corporate": {
        "accent": "#dc2626", "dark": "#1e293b", "muted": "#78716c",
        "light": "#fef2f2", "background": "#ffffff",
        "secondary": "#0f172a", "success": "#15803d", "warning": "#ca8a04",
        "font_title": "Arial", "font_body": "Arial",
    },
}

# ═══════════════════════════════════════════════════════════════
#  SLIDE LAYOUTS — expanded catalog
# ═══════════════════════════════════════════════════════════════

LAYOUTS = {
    "title_slide": "Title + subtitle + author",
    "toc": "Agenda / table of contents",
    "section_header": "Section divider with large title",
    "content_bullets": "Title + bullet points",
    "content_body": "Title + freeform body text",
    "two_column": "Title + 2-column layout (text | text)",
    "image_right": "Title + bullets left, image right",
    "image_left": "Title + image left, bullets right",
    "image_full": "Full-bleed image with overlay title",
    "table": "Title + data table",
    "chart": "Title + chart (bar/line/pie)",
    "chart_with_text": "Title + chart left, insights right",
    "quote": "Large quote + attribution",
    "comparison": "Title + 2-column comparison (A vs B)",
    "timeline": "Title + horizontal timeline",
    "kpi_dashboard": "3-4 KPI cards with numbers",
    "team": "Team member cards with photos",
    "thank_you": "Closing slide with contact info",
    "blank": "Blank canvas",
}

TEMPLATE_CATALOG = [
    {"name": "professional", "category": "business", "description": "Clean business reporting."},
    {"name": "academic", "category": "education", "description": "Research, thesis, lecture."},
    {"name": "creative", "category": "creative", "description": "Expressive proposals, marketing."},
    {"name": "minimal", "category": "simple", "description": "Sparse internal communication."},
    {"name": "dark", "category": "modern", "description": "Dark theme for modern presentations."},
    {"name": "corporate", "category": "business", "description": "Bold corporate branding."},
]

# ═══════════════════════════════════════════════════════════════
#  UTILITIES
# ═══════════════════════════════════════════════════════════════

SOURCE_CONVERTERS = {
    ".pdf": ["scripts/source_to_md/pdf_to_md.py"],
    ".docx": ["scripts/source_to_md/doc_to_md.py"],
    ".doc": ["scripts/source_to_md/doc_to_md.py"],
    ".xlsx": ["scripts/source_to_md/excel_to_md.py"],
    ".xlsm": ["scripts/source_to_md/excel_to_md.py"],
    ".xls": ["scripts/source_to_md/excel_to_md.py"],
    ".pptx": ["scripts/source_to_md/ppt_to_md.py"],
    ".ppt": ["scripts/source_to_md/ppt_to_md.py"],
    ".html": ["scripts/source_to_md/html_to_md.py"],
    ".htm": ["scripts/source_to_md/html_to_md.py"],
    ".epub": ["scripts/source_to_md/epub_to_md.py"],
    ".tex": ["scripts/source_to_md/latex_to_md.py"],
    ".rst": ["scripts/source_to_md/rst_to_md.py"],
}

UPSTREAM_REPO_URL = "https://github.com/hugohe3/ppt-master.git"
DEFAULT_EXTERNAL_ROOT = Path.home() / ".openmegatron" / "ppt-master"
REQUIRED_EXTERNAL_SCRIPTS = [
    "scripts/finalize_svg.py",
    "scripts/svg_to_pptx.py",
]
OPTIONAL_EXTERNAL_IMPORTS = [
    ("pptx", "python-pptx"),
    ("fitz", "PyMuPDF"),
]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _hex_to_rgb(value: str) -> RGBColor:
    text = str(value or "").strip().lstrip("#")
    if len(text) == 3:
        text = "".join(ch * 2 for ch in text)
    if not re.fullmatch(r"[0-9a-fA-F]{6}", text):
        text = "2563eb"
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def _safe_filename(value: str) -> str:
    cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", str(value or "deck")).strip(" ._")
    return cleaned[:80] or "deck"


def _text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, (list, tuple)):
        return "\n".join(str(item) for item in value)
    return str(value)


def _as_bullets(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, str):
        return [re.sub(r"^\s*[-*+]\s+", "", line).strip()
                for line in value.splitlines() if line.strip()]
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()]
    return [str(value).strip()] if str(value).strip() else []


def _truthy(value: Any, default: bool = False) -> bool:
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    return str(value).strip().lower() in {"1", "true", "yes", "y", "on"}


def _default_output_path(args: dict, title: str) -> Path:
    output = args.get("output")
    if output:
        out_path = Path(str(output)).expanduser()
    else:
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        out_path = Path("outputs") / "ppt_master" / f"{_safe_filename(title)}_{stamp}.pptx"
    if out_path.suffix.lower() != ".pptx":
        out_path = out_path.with_suffix(".pptx")
    return out_path


def _ensure_output_writable(path: Path, overwrite: bool) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    if path.exists() and not overwrite:
        raise FileExistsError(f"Output exists: {path}. Pass overwrite=true to replace it.")


def _is_url(value: str) -> bool:
    return bool(re.match(r"^https?://", str(value or ""), re.I))


def _run_command(command: list[str], cwd: Path | None, timeout_seconds: int) -> subprocess.CompletedProcess:
    return subprocess.run(
        [str(part) for part in command],
        cwd=str(cwd) if cwd else None,
        text=True, encoding="utf-8", errors="replace",
        capture_output=True, timeout=timeout_seconds,
    )


def _python_executable(args: dict) -> str:
    return str(args.get("python_executable") or sys.executable)


def parse_cli_args() -> dict:
    if len(sys.argv) <= 1:
        return {}
    raw = sys.argv[1]
    try:
        parsed = json.loads(raw)
        if isinstance(parsed, dict):
            return parsed
    except Exception:
        pass
    return {}


# ═══════════════════════════════════════════════════════════════
#  BUILTIN RENDERER — upgraded with layouts, charts, tables
# ═══════════════════════════════════════════════════════════════


def _make_font(size: int, color: str, bold: bool = False, italic: bool = False,
               name: str = None) -> dict:
    return {"size": Pt(size), "color": color, "bold": bold, "italic": italic, "name": name}


def _apply_font(run, font_spec: dict, theme: dict) -> None:
    run.font.size = font_spec.get("size", Pt(16))
    run.font.bold = font_spec.get("bold", False)
    if font_spec.get("italic"):
        run.font.italic = True
    run.font.color.rgb = _hex_to_rgb(font_spec.get("color", theme["dark"]))
    name = font_spec.get("name") or theme.get("font_body")
    if name:
        run.font.name = name


def _add_textbox(slide, left, top, width, height, text: str,
                 font_spec: dict, theme: dict, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, word_wrap: bool = True,
                 line_spacing: float = 1.15) -> Any:
    box = slide.shapes.add_textbox(
        Inches(left) if isinstance(left, (int, float)) else left,
        Inches(top) if isinstance(top, (int, float)) else top,
        Inches(width) if isinstance(width, (int, float)) else width,
        Inches(height) if isinstance(height, (int, float)) else height,
    )
    frame = box.text_frame
    frame.word_wrap = word_wrap
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.alignment = align
    if line_spacing != 1.0:
        try:
            pts = Pt(int(Pt(18) / 12700 * line_spacing * 12700))
            if int(pts) > 20116800:
                pts = Pt(20116800)
            paragraph.line_spacing = pts
        except (ValueError, TypeError):
            pass  # skip line spacing if invalid
    if anchor != MSO_ANCHOR.TOP:
        frame.paragraphs[0].space_before = Pt(0)
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    _apply_font(run, font_spec, theme)
    return box


def _add_bullets(slide, bullets: list[str], left, top, width, height,
                 font_spec: dict, theme: dict, line_spacing: float = 1.3) -> Any:
    box = slide.shapes.add_textbox(
        Inches(left) if isinstance(left, (int, float)) else left,
        Inches(top) if isinstance(top, (int, float)) else top,
        Inches(width) if isinstance(width, (int, float)) else width,
        Inches(height) if isinstance(height, (int, float)) else height,
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.08)
    for i, bullet in enumerate(bullets):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.text = bullet
        para.level = 0
        para.space_after = Pt(6)
        run = para.runs[0] if para.runs else para.add_run()
        _apply_font(run, font_spec, theme)
    return box


def _add_shape_bg(slide, left, top, width, height, color: str) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left) if isinstance(left, (int, float)) else left,
        Inches(top) if isinstance(top, (int, float)) else top,
        Inches(width) if isinstance(width, (int, float)) else width,
        Inches(height) if isinstance(height, (int, float)) else height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(color)
    shape.line.fill.background()  # no border
    return shape


def _add_table(slide, table_data: Any, left, top, width, height, theme: dict) -> bool:
    """Render a styled data table. Supports dict-of-lists, list-of-dicts, list-of-lists."""
    if not isinstance(table_data, (list, dict)):
        return False
    rows: list[list[str]] = []
    if isinstance(table_data, dict):
        headers = list(table_data.keys())
        max_len = max(len(v) if isinstance(v, list) else 1 for v in table_data.values())
        rows.append(headers)
        for i in range(max_len):
            row = []
            for h in headers:
                val = table_data.get(h, [])
                if isinstance(val, list) and i < len(val):
                    row.append(str(val[i]))
                elif i == 0 and not isinstance(val, list):
                    row.append(str(val))
                else:
                    row.append("")
            rows.append(row)
    elif all(isinstance(row, dict) for row in table_data):
        headers = list(table_data[0].keys())
        rows.append(headers)
        rows.extend([[str(row.get(h, "")) for h in headers] for row in table_data])
    elif all(isinstance(row, list) for row in table_data):
        rows = [[str(cell) for cell in row] for row in table_data]
    if not rows:
        return False

    max_cols = min(max(len(row) for row in rows), 7)
    max_rows = min(len(rows), 12)
    shape = slide.shapes.add_table(
        max_rows, max_cols,
        Inches(left) if isinstance(left, (int, float)) else left,
        Inches(top) if isinstance(top, (int, float)) else top,
        Inches(width) if isinstance(width, (int, float)) else width,
        Inches(height) if isinstance(height, (int, float)) else height,
    )
    tbl = shape.table

    # Set column widths evenly
    col_w = int(Inches(width) / max_cols) if isinstance(width, (int, float)) else int(Emu(12192000) / max_cols)
    for c in range(max_cols):
        tbl.columns[c].width = col_w

    for r in range(max_rows):
        for c in range(max_cols):
            cell = tbl.cell(r, c)
            cell.text = rows[r][c] if c < len(rows[r]) else ""
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.color.rgb = _hex_to_rgb(theme["dark"])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Header row styling
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex_to_rgb(theme["accent"])
                para.font.bold = True
                para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                para.font.size = Pt(12)
            elif r % 2 == 0:
                # Alternating row color
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex_to_rgb(theme["light"])
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # Thin border
            tcPr = cell._tc.get_or_add_tcPr()
            for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = tcPr.find(qn(border_name))
                if ln is None:
                    from lxml import etree
                    ln = etree.SubElement(tcPr, qn(border_name))
                ln.set('w', '6350')  # 0.5pt
                solidFill = ln.find(qn('a:solidFill'))
                if solidFill is None:
                    from lxml import etree
                    solidFill = etree.SubElement(ln, qn('a:solidFill'))
                srgb = solidFill.find(qn('a:srgbClr'))
                if srgb is None:
                    from lxml import etree
                    srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
                srgb.set('val', 'd1d5db')

    return True


def _add_chart(slide, chart_data: dict, left, top, width, height, theme: dict) -> bool:
    """Add a native PowerPoint chart (bar, line, or pie).

    chart_data format:
      {"type": "bar"|"line"|"pie", "categories": [...], "series": [{"name": "...", "values": [...]}]}
    """
    if not isinstance(chart_data, dict):
        return False

    chart_type_name = str(chart_data.get("type", "bar")).lower()
    chart_type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
        "bar_horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
    }
    chart_type = chart_type_map.get(chart_type_name, XL_CHART_TYPE.COLUMN_CLUSTERED)

    categories = chart_data.get("categories", [])
    series_list = chart_data.get("series", [])
    if not categories or not series_list:
        return False

    chart_frame = slide.shapes.add_chart(
        chart_type,
        Inches(left) if isinstance(left, (int, float)) else left,
        Inches(top) if isinstance(top, (int, float)) else top,
        Inches(width) if isinstance(width, (int, float)) else width,
        Inches(height) if isinstance(height, (int, float)) else height,
    )
    chart = chart_frame.chart

    # Replace placeholder data
    plot = chart.plots[0]
    for i, series_info in enumerate(series_list):
        values = series_info.get("values", [])
        name = series_info.get("name", f"Series {i + 1}")
        if i == 0:
            # Use first series as the primary
            chart_data_obj = chart_frame.chart.series[0] if chart.series else None
        # Build category axis
        cat_axis = chart.category_axis if hasattr(chart, 'category_axis') else None

    # Rebuild chart data properly
    from pptx.chart.data import CategoryChartData
    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = categories
    for series_info in series_list:
        chart_data_obj.add_series(
            series_info.get("name", "Series"),
            series_info.get("values", []),
        )

    # Replace the chart's data
    chart.replace_data(chart_data_obj)

    # Style the chart
    chart.has_legend = len(series_list) > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)

    # Color the series
    accent_colors = [theme["accent"], theme["secondary"], theme["success"],
                     theme["warning"], theme["muted"]]
    for i, series in enumerate(chart.series):
        color = accent_colors[i % len(accent_colors)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _hex_to_rgb(color)

    return True


# ═══════════════════════════════════════════════════════════════
#  SLIDE BUILDERS — one function per layout type
# ═══════════════════════════════════════════════════════════════


def _build_title_slide(prs: Presentation, title: str, subtitle: str,
                       author: str, date_str: str, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _hex_to_rgb(theme["background"])

    # Full-width accent bar at top
    _add_shape_bg(slide, 0, 0, 13.333, 0.15, theme["accent"])

    # Title
    _add_textbox(slide, 1.0, 2.0, 11.3, 1.4, title,
                 _make_font(42, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    # Subtitle
    if subtitle:
        _add_textbox(slide, 1.05, 3.6, 10.3, 0.7, subtitle,
                     _make_font(20, theme["muted"], name=theme.get("font_body")),
                     theme)

    # Author + date
    meta_parts = []
    if author:
        meta_parts.append(author)
    if date_str:
        meta_parts.append(date_str)
    if meta_parts:
        _add_textbox(slide, 1.05, 6.5, 10.3, 0.35, "  |  ".join(meta_parts),
                     _make_font(12, theme["muted"], name=theme.get("font_body")),
                     theme)

    # Bottom accent line
    _add_shape_bg(slide, 0, 7.35, 13.333, 0.04, theme["accent"])


def _build_toc_slide(prs: Presentation, slides_data: list[dict], theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_bg(slide, 0, 0, 13.333, 0.08, theme["accent"])
    _add_textbox(slide, 0.8, 0.5, 11.5, 0.7, "Contents",
                 _make_font(30, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    items = []
    for i, s in enumerate(slides_data):
        num = f"{i + 1:02d}"
        title = str(s.get("title") or f"Slide {i + 1}")
        items.append(f"{num}    {title}")

    _add_bullets(slide, items, 1.2, 1.5, 10.8, 5.2,
                 _make_font(18, theme["dark"], name=theme.get("font_body")),
                 theme, line_spacing=1.6)


def _build_section_header(prs: Presentation, title: str, subtitle: str,
                          section_num: int, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_bg(slide, 0, 0, 13.333, 7.5, theme["accent"])
    if section_num:
        _add_textbox(slide, 1.0, 2.0, 11.3, 0.6, f"PART {section_num:02d}",
                     _make_font(14, "#ffffffcc", bold=True, name=theme.get("font_body")),
                     theme)
    _add_textbox(slide, 1.0, 2.7, 11.3, 1.5, title,
                 _make_font(40, "#ffffff", bold=True, name=theme.get("font_title")),
                 theme)
    if subtitle:
        _add_textbox(slide, 1.05, 4.4, 10.3, 0.6, subtitle,
                     _make_font(18, "#ffffffcc", name=theme.get("font_body")),
                     theme)


def _build_content_bullets(prs: Presentation, item: dict, index: int,
                           theme: dict, image_placeholders: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_bg(slide, 0.6, 0.35, 0.06, 0.55, theme["accent"])

    title = str(item.get("title") or f"Slide {index}")
    _add_textbox(slide, 0.78, 0.32, 11.5, 0.7, title,
                 _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    subtitle = _text(item.get("subtitle"))
    if subtitle:
        _add_textbox(slide, 0.8, 0.95, 11.0, 0.35, subtitle,
                     _make_font(13, theme["muted"], name=theme.get("font_body")),
                     theme)

    bullets = _as_bullets(item.get("bullets") or item.get("points") or item.get("body"))
    if bullets:
        size = 22 if len(bullets) <= 4 else 18
        _add_bullets(slide, bullets, 1.1, 1.6, 11.0, 5.0,
                     _make_font(size, theme["dark"], name=theme.get("font_body")),
                     theme)

    # Page number
    _add_textbox(slide, 12.5, 7.1, 0.6, 0.2, str(index),
                 _make_font(9, theme["muted"]), theme, align=PP_ALIGN.RIGHT)


def _build_two_column(prs: Presentation, item: dict, index: int, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_bg(slide, 0.6, 0.35, 0.06, 0.55, theme["accent"])
    title = str(item.get("title") or f"Slide {index}")
    _add_textbox(slide, 0.78, 0.32, 11.5, 0.7, title,
                 _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    left_bullets = _as_bullets(item.get("left") or item.get("column_a") or [])
    right_bullets = _as_bullets(item.get("right") or item.get("column_b") or [])
    if not left_bullets and not right_bullets:
        all_bullets = _as_bullets(item.get("bullets", []))
        mid = len(all_bullets) // 2
        left_bullets = all_bullets[:mid]
        right_bullets = all_bullets[mid:]

    if left_bullets:
        _add_bullets(slide, left_bullets, 0.8, 1.5, 5.5, 5.2,
                     _make_font(18, theme["dark"], name=theme.get("font_body")), theme)
    if right_bullets:
        _add_bullets(slide, right_bullets, 7.0, 1.5, 5.5, 5.2,
                     _make_font(18, theme["dark"], name=theme.get("font_body")), theme)


def _build_image_with_text(prs: Presentation, item: dict, index: int,
                           theme: dict, image_placeholders: bool, image_left: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_bg(slide, 0.6, 0.35, 0.06, 0.55, theme["accent"])
    title = str(item.get("title") or f"Slide {index}")
    _add_textbox(slide, 0.78, 0.32, 11.5, 0.7, title,
                 _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    image_path = _text(item.get("image") or item.get("image_path"))
    bullets = _as_bullets(item.get("bullets") or item.get("points"))

    text_left = 0.8 if image_left else 6.6
    img_left = 6.6 if image_left else 0.8

    if bullets:
        _add_bullets(slide, bullets[:6], text_left, 1.5, 5.4, 5.0,
                     _make_font(18, theme["dark"], name=theme.get("font_body")), theme)

    if image_path and Path(image_path).expanduser().exists():
        slide.shapes.add_picture(
            str(Path(image_path).expanduser()),
            Inches(img_left), Inches(1.5), width=Inches(5.6), height=Inches(4.8),
        )
    elif image_placeholders:
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(img_left), Inches(1.5), Inches(5.6), Inches(4.8),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex_to_rgb(theme["light"])
        rect.line.color.rgb = _hex_to_rgb(theme["muted"])
        rect.line.width = Pt(0.5)
        _add_textbox(slide, img_left + 0.3, 3.6, 5.0, 0.4,
                     "Image placeholder",
                     _make_font(14, theme["muted"], name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER)


def _build_quote_slide(prs: Presentation, item: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_bg(slide, 0, 0, 13.333, 7.5, theme["accent"])

    quote = str(item.get("quote") or item.get("body") or "")
    if not quote.startswith('"'):
        quote = f'"{quote}"'

    _add_textbox(slide, 1.5, 2.0, 10.3, 3.0, quote,
                 _make_font(32, "#ffffff", italic=True, name=theme.get("font_title")),
                 theme, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    attribution = str(item.get("attribution") or item.get("author") or "")
    if attribution:
        _add_textbox(slide, 1.5, 5.5, 10.3, 0.5, f"— {attribution}",
                     _make_font(16, "#ffffffcc", name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER)


def _build_comparison_slide(prs: Presentation, item: dict, index: int, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_bg(slide, 0.6, 0.35, 0.06, 0.55, theme["accent"])
    title = str(item.get("title") or f"Slide {index}")
    _add_textbox(slide, 0.78, 0.32, 11.5, 0.7, title,
                 _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    left_title = str(item.get("left_title") or "Option A")
    right_title = str(item.get("right_title") or "Option B")
    _add_textbox(slide, 1.0, 1.3, 5.2, 0.5, left_title,
                 _make_font(20, theme["accent"], bold=True, name=theme.get("font_body")),
                 theme, align=PP_ALIGN.CENTER)
    _add_textbox(slide, 7.1, 1.3, 5.2, 0.5, right_title,
                 _make_font(20, theme["secondary"], bold=True, name=theme.get("font_body")),
                 theme, align=PP_ALIGN.CENTER)

    # Divider
    _add_shape_bg(slide, 6.55, 1.3, 0.03, 5.2, theme["muted"])

    left_items = _as_bullets(item.get("left_items") or [])
    right_items = _as_bullets(item.get("right_items") or [])
    if left_items:
        _add_bullets(slide, left_items, 0.6, 2.0, 5.6, 4.5,
                     _make_font(16, theme["dark"], name=theme.get("font_body")), theme)
    if right_items:
        _add_bullets(slide, right_items, 7.0, 2.0, 5.6, 4.5,
                     _make_font(16, theme["dark"], name=theme.get("font_body")), theme)


def _build_timeline_slide(prs: Presentation, item: dict, index: int, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_bg(slide, 0.6, 0.35, 0.06, 0.55, theme["accent"])
    title = str(item.get("title") or f"Slide {index}")
    _add_textbox(slide, 0.78, 0.32, 11.5, 0.7, title,
                 _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    milestones = item.get("milestones") or item.get("events") or []
    if not isinstance(milestones, list):
        return

    # Horizontal line
    _add_shape_bg(slide, 0.8, 3.7, 11.7, 0.04, theme["accent"])

    n = len(milestones)
    spacing = 11.5 / max(n, 1)
    for i, m in enumerate(milestones):
        x = 0.8 + i * spacing + spacing / 2
        # Dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x - 0.15), Inches(3.55), Inches(0.3), Inches(0.3),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = _hex_to_rgb(theme["accent"])
        dot.line.fill.background()

        # Label
        label = str(m.get("label") or m.get("title") or f"Event {i + 1}")
        date_str = str(m.get("date") or m.get("time") or "")

        # Alternate above/below the line
        y_offset = -1.0 if i % 2 == 0 else 0.5
        _add_textbox(slide, x - 1.5, 3.7 + y_offset, 3.0, 0.4, date_str,
                     _make_font(11, theme["accent"], bold=True, name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER)
        _add_textbox(slide, x - 1.5, 3.7 + y_offset + 0.35, 3.0, 0.6, label,
                     _make_font(14, theme["dark"], name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER)


def _build_kpi_dashboard(prs: Presentation, item: dict, index: int, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_bg(slide, 0.6, 0.35, 0.06, 0.55, theme["accent"])
    title = str(item.get("title") or f"Slide {index}")
    _add_textbox(slide, 0.78, 0.32, 11.5, 0.7, title,
                 _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    kpis = item.get("kpis") or item.get("metrics") or []
    if not isinstance(kpis, list):
        return

    n = min(len(kpis), 4)
    card_w = 11.5 / n - 0.3
    for i, kpi in enumerate(kpis[:4]):
        x = 0.8 + i * (card_w + 0.3)
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5),
            Inches(card_w), Inches(2.8),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _hex_to_rgb(theme["light"])
        card.line.fill.background()

        value = str(kpi.get("value") or kpi.get("number") or "")
        label = str(kpi.get("label") or kpi.get("title") or "")
        change = str(kpi.get("change") or kpi.get("trend") or "")

        _add_textbox(slide, x + 0.2, 1.8, card_w - 0.4, 1.2, value,
                     _make_font(36, theme["accent"], bold=True, name=theme.get("font_title")),
                     theme, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_textbox(slide, x + 0.2, 3.0, card_w - 0.4, 0.5, label,
                     _make_font(14, theme["dark"], name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER)
        if change:
            change_color = theme["success"] if change.startswith("+") else theme["warning"] if change.startswith("-") else theme["muted"]
            _add_textbox(slide, x + 0.2, 3.5, card_w - 0.4, 0.4, change,
                         _make_font(13, change_color, bold=True, name=theme.get("font_body")),
                         theme, align=PP_ALIGN.CENTER)


def _build_thank_you(prs: Presentation, item: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_bg(slide, 0, 0, 13.333, 7.5, theme["accent"])

    text = str(item.get("title") or "Thank You")
    _add_textbox(slide, 1.5, 2.5, 10.3, 1.5, text,
                 _make_font(48, "#ffffff", bold=True, name=theme.get("font_title")),
                 theme, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    contact = str(item.get("contact") or item.get("subtitle") or "")
    if contact:
        _add_textbox(slide, 1.5, 4.5, 10.3, 0.8, contact,
                     _make_font(18, "#ffffffcc", name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER)


def _build_full_image_slide(prs: Presentation, item: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_path = _text(item.get("image") or item.get("image_path"))
    title = str(item.get("title") or "")

    if image_path and Path(image_path).expanduser().exists():
        slide.shapes.add_picture(
            str(Path(image_path).expanduser()),
            Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5),
        )

    if title:
        # Dark overlay bar at bottom
        _add_shape_bg(slide, 0, 5.8, 13.333, 1.7, "#00000088")
        _add_textbox(slide, 0.8, 6.0, 11.5, 0.8, title,
                     _make_font(28, "#ffffff", bold=True, name=theme.get("font_title")),
                     theme)


def _build_team_slide(prs: Presentation, item: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape_bg(slide, 0.6, 0.35, 0.06, 0.55, theme["accent"])
    title = str(item.get("title") or "Team")
    _add_textbox(slide, 0.78, 0.32, 11.5, 0.7, title,
                 _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                 theme)

    members = item.get("members") or item.get("team") or []
    if not isinstance(members, list):
        return

    n = min(len(members), 5)
    card_w = 11.5 / n - 0.4
    for i, member in enumerate(members[:5]):
        x = 0.8 + i * (card_w + 0.4)
        name = str(member.get("name") or "")
        role = str(member.get("role") or member.get("title") or "")

        # Avatar placeholder circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + card_w / 2 - 0.7), Inches(1.6),
            Inches(1.4), Inches(1.4),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = _hex_to_rgb(theme["light"])
        circle.line.color.rgb = _hex_to_rgb(theme["accent"])
        circle.line.width = Pt(2)

        # Initials
        initials = "".join(w[0].upper() for w in name.split()[:2]) if name else "?"
        _add_textbox(slide, x + card_w / 2 - 0.7, 1.9, 1.4, 0.6, initials,
                     _make_font(24, theme["accent"], bold=True, name=theme.get("font_title")),
                     theme, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        _add_textbox(slide, x, 3.3, card_w, 0.5, name,
                     _make_font(16, theme["dark"], bold=True, name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER)
        _add_textbox(slide, x, 3.75, card_w, 0.4, role,
                     _make_font(13, theme["muted"], name=theme.get("font_body")),
                     theme, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
#  SLIDE DISPATCHER
# ═══════════════════════════════════════════════════════════════

def _build_content_slide(prs: Presentation, item: dict, index: int,
                         theme: dict, image_placeholders: bool) -> None:
    """Route to the appropriate slide builder based on layout field."""
    layout = str(item.get("layout") or "").lower()

    if layout in ("title_slide", "title"):
        _build_title_slide(prs, str(item.get("title") or ""),
                          str(item.get("subtitle") or ""),
                          str(item.get("author") or ""),
                          str(item.get("date") or ""), theme)
    elif layout == "section_header":
        _build_section_header(prs, str(item.get("title") or ""),
                             str(item.get("subtitle") or ""),
                             int(item.get("section_num") or 0), theme)
    elif layout == "two_column":
        _build_two_column(prs, item, index, theme)
    elif layout == "image_right":
        _build_image_with_text(prs, item, index, theme, image_placeholders, image_left=False)
    elif layout == "image_left":
        _build_image_with_text(prs, item, index, theme, image_placeholders, image_left=True)
    elif layout == "image_full":
        _build_full_image_slide(prs, item, theme)
    elif layout == "quote":
        _build_quote_slide(prs, item, theme)
    elif layout == "comparison":
        _build_comparison_slide(prs, item, index, theme)
    elif layout == "timeline":
        _build_timeline_slide(prs, item, index, theme)
    elif layout == "kpi_dashboard":
        _build_kpi_dashboard(prs, item, index, theme)
    elif layout == "team":
        _build_team_slide(prs, item, theme)
    elif layout == "thank_you":
        _build_thank_you(prs, item, theme)
    elif layout == "chart" or layout == "chart_with_text":
        _build_content_bullets(prs, item, index, theme, image_placeholders)
        # Add chart on a separate slide for now
        chart_data = item.get("chart")
        if chart_data:
            chart_slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_shape_bg(chart_slide, 0.6, 0.35, 0.06, 0.55, theme["accent"])
            chart_title = str(item.get("chart_title") or item.get("title") or "Chart")
            _add_textbox(chart_slide, 0.78, 0.32, 11.5, 0.7, chart_title,
                         _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                         theme)
            _add_chart(chart_slide, chart_data, 0.8, 1.5, 11.5, 5.2, theme)
    elif layout == "table":
        _build_content_bullets(prs, item, index, theme, image_placeholders)
        table_data = item.get("table")
        if table_data:
            table_slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_shape_bg(table_slide, 0.6, 0.35, 0.06, 0.55, theme["accent"])
            table_title = str(item.get("table_title") or item.get("title") or "Data")
            _add_textbox(table_slide, 0.78, 0.32, 11.5, 0.7, table_title,
                         _make_font(28, theme["dark"], bold=True, name=theme.get("font_title")),
                         theme)
            _add_table(table_slide, table_data, 0.8, 1.5, 11.5, 5.2, theme)
    else:
        # Default: bullets layout
        _build_content_bullets(prs, item, index, theme, image_placeholders)


# ═══════════════════════════════════════════════════════════════
#  MARKDOWN / OUTLINE PARSING
# ═══════════════════════════════════════════════════════════════

def _parse_markdown_outline(outline: str) -> list[dict]:
    slides: list[dict] = []
    current: dict | None = None
    for raw_line in str(outline or "").splitlines():
        line = raw_line.rstrip()
        # H1 = section header, H2 = slide title, H3 = layout hint
        h1 = re.match(r"^\s*#\s+(.+?)\s*$", line)
        h2 = re.match(r"^\s*#{2}\s+(.+?)\s*$", line)
        h3 = re.match(r"^\s*#{3}\s+(.+?)\s*$", line)

        if h1:
            if current:
                slides.append(current)
            current = {"title": h1.group(1).strip(), "layout": "section_header", "bullets": []}
        elif h2:
            if current:
                slides.append(current)
            current = {"title": h2.group(1).strip(), "bullets": []}
        elif h3:
            # Layout hint: ## Slide title {layout=comparison}
            layout_match = re.search(r"\{layout=(\w+)\}", h3.group(1))
            if current:
                if layout_match:
                    current["layout"] = layout_match.group(1)
            elif layout_match:
                current = {"title": h3.group(1).strip(), "layout": layout_match.group(1), "bullets": []}
        else:
            bullet = re.match(r"^\s*[-*+]\s+(.+?)\s*$", line)
            table_row = re.match(r"^\s*\|(.+)\|\s*$", line)
            if bullet:
                if current is None:
                    current = {"title": "Overview", "bullets": []}
                current.setdefault("bullets", []).append(bullet.group(1).strip())
            elif table_row and current:
                cells = [c.strip() for c in table_row.group(1).split("|")]
                current.setdefault("table_rows", []).append(cells)
            elif line.strip():
                if current is None:
                    current = {"title": line.strip(), "bullets": []}
                else:
                    current.setdefault("body", "")
                    current["body"] = (current["body"] + "\n" + line.strip()).strip()

    if current:
        # Convert table_rows to proper table format if present
        if current.get("table_rows"):
            rows = current.pop("table_rows")
            if len(rows) > 1:
                current["table"] = [dict(zip(rows[0], row)) for row in rows[1:]]
        slides.append(current)
    return slides


def _fallback_slides(topic: str, audience: str, language: str, slide_count: int) -> list[dict]:
    zh = str(language or "").lower().startswith("zh") or bool(re.search(r"[一-鿿]", topic or ""))
    target = audience or ("业务与技术团队" if zh else "business and technical audience")
    if zh:
        templates = [
            {"title": "核心结论", "bullets": [f"{topic} 的关键价值与适用边界", "用一句话说明问题、方案和预期收益", f"面向 {target} 保持表达清晰可执行"]},
            {"title": "背景与痛点", "bullets": ["现有流程中的低效、重复和信息断层", "用户真正需要的是更快生成、可持续编辑和可复用模板", "明确成功标准和约束条件"]},
            {"title": "方案架构", "layout": "two_column",
             "left": ["输入：主题、资料、风格、页数与受众", "处理：大纲生成、版式规划、内容填充", "输出：原生可编辑 PPTX 文件"],
             "right": ["文本排版、表格渲染、图表生成", "智能配色、布局自适应", "支持 PDF/DOCX 源文件转换"]},
            {"title": "关键指标", "layout": "kpi_dashboard",
             "kpis": [
                 {"value": "6", "label": "主题风格", "change": "+2"},
                 {"value": "16", "label": "版式布局", "change": "+10"},
                 {"value": "100%", "label": "可编辑性", "change": "+100%"},
                 {"value": "<3s", "label": "平均生成时间", "change": "-80%"},
             ]},
            {"title": "落地流程", "layout": "timeline",
             "milestones": [
                 {"date": "Day 1", "label": "收集资料并生成初稿"},
                 {"date": "Day 2", "label": "人工校对事实与叙事节奏"},
                 {"date": "Day 3", "label": "统一视觉主题"},
                 {"date": "Day 4", "label": "导出交付"},
             ]},
            {"title": "下一步", "bullets": ["补充真实数据和品牌模板", "接入图表、图片和引用来源", "沉淀为可复用的演示文稿工作流"]},
        ]
    else:
        templates = [
            {"title": "Executive Summary", "bullets": [f"Key value and boundaries for {topic}", "State the problem, solution, and expected outcome", f"Keep the deck actionable for {target}"]},
            {"title": "Context and Pain Points", "bullets": ["Current workflows are slow, repetitive, or hard to reuse", "Users need editable outputs, not static screenshots", "Define success criteria and constraints early"]},
            {"title": "Solution Architecture", "layout": "two_column",
             "left": ["Input: topic, source material, style, audience", "Process: outline, layout, content, placement", "Output: native editable PPTX"],
             "right": ["Text formatting, table rendering, charts", "Smart color palettes, auto-layout", "Source conversion from PDF/DOCX"]},
            {"title": "Key Metrics", "layout": "kpi_dashboard",
             "kpis": [
                 {"value": "6", "label": "Theme Styles", "change": "+2"},
                 {"value": "16", "label": "Layout Types", "change": "+10"},
                 {"value": "100%", "label": "Editable", "change": "+100%"},
                 {"value": "<3s", "label": "Avg Generation", "change": "-80%"},
             ]},
            {"title": "Workflow", "layout": "timeline",
             "milestones": [
                 {"date": "Day 1", "label": "Gather materials, generate draft"},
                 {"date": "Day 2", "label": "Review facts and narrative"},
                 {"date": "Day 3", "label": "Apply visual theme"},
                 {"date": "Day 4", "label": "Export and deliver"},
             ]},
            {"title": "Next Steps", "bullets": ["Add real data and brand templates", "Connect charts, images, and citations", "Turn into reusable presentation workflow"]},
        ]
    count = max(1, min(int(slide_count or 6), len(templates)))
    return templates[:count]


def _normalize_slides(args: dict) -> list[dict]:
    slides = args.get("slides")
    if isinstance(slides, str):
        try:
            slides = json.loads(slides)
        except Exception:
            slides = None
    if isinstance(slides, list) and slides:
        return slides

    outline = args.get("outline")
    if outline:
        parsed = _parse_markdown_outline(str(outline))
        if parsed:
            return parsed

    return _fallback_slides(
        topic=str(args.get("topic") or args.get("title") or "Presentation"),
        audience=str(args.get("audience") or ""),
        language=str(args.get("language") or "zh"),
        slide_count=int(args.get("slide_count") or 6),
    )


# ═══════════════════════════════════════════════════════════════
#  MAIN BUILDER
# ═══════════════════════════════════════════════════════════════

def build_presentation(args: dict, engine: str = "builtin") -> dict:
    topic = str(args.get("topic") or args.get("title") or "Presentation").strip()
    title = str(args.get("title") or topic).strip()
    style = str(args.get("style") or "professional").strip().lower()
    theme = dict(THEMES.get(style, THEMES["professional"]))
    if args.get("theme_color"):
        theme["accent"] = str(args["theme_color"])

    # Load brand config
    brand = _load_brand_config(args)
    if brand:
        if brand.get("primary"):
            theme["accent"] = brand["primary"]
        if brand.get("font_title"):
            theme["font_title"] = brand["font_title"]
        if brand.get("font_body"):
            theme["font_body"] = brand["font_body"]

    # AI content generation
    slides_data = _normalize_slides(args)
    if args.get("ai_generate") and not args.get("slides") and not args.get("outline"):
        ai_slides = _ai_generate_slides(args, theme, brand)
        if ai_slides:
            slides_data = ai_slides

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    image_placeholders = bool(args.get("image_placeholders", True))

    # Apply brand logo to slide master
    if brand:
        _apply_slide_master(prs, brand, theme)

    # Title slide
    _build_title_slide(prs, title=title,
                       subtitle=str(args.get("subtitle") or ""),
                       author=str(args.get("author") or (brand.get("company", "") if brand else "")),
                       date_str=str(args.get("date") or datetime.now().strftime("%Y-%m-%d")),
                       theme=theme)

    # TOC
    if args.get("include_toc", True) and len(slides_data) > 1:
        _build_toc_slide(prs, slides_data, theme)

    # Content slides with brand footer
    for index, item in enumerate(slides_data, 1):
        _build_content_slide(prs, item, index, theme, image_placeholders)
        if brand:
            _apply_brand_footer(prs.slides[-1], brand, theme, index, len(slides_data))

    # Thank you slide
    include_thanks = args.get("include_thanks", True)
    if include_thanks:
        _build_thank_you(prs, {
            "title": str(args.get("closing_text") or "Thank You"),
            "subtitle": str(args.get("closing_subtitle") or args.get("contact") or
                          (brand.get("footer", "") if brand else "")),
        }, theme)

    out_path = _default_output_path(args, title)
    _ensure_output_writable(out_path, _truthy(args.get("overwrite"), False))
    prs.save(out_path)

    return {
        "status": "success",
        "format": "pptx",
        "output": str(out_path.resolve()),
        "slide_count": len(prs.slides),
        "content_slides": len(slides_data),
        "editable": True,
        "engine": engine,
        "style": style,
        "theme": style,
        "ai_generated": bool(args.get("ai_generate") and not args.get("slides") and not args.get("outline")),
        "brand_applied": brand is not None,
    }


# ═══════════════════════════════════════════════════════════════
#  BRAND TEMPLATE SYSTEM
# ═══════════════════════════════════════════════════════════════

DEFAULT_BRAND_PATH = Path("pysrc/skills/office/ppt-master-1.0.0/brand.toml")


def _load_brand_config(args: dict) -> dict | None:
    """Load brand configuration from args or brand.toml file."""
    brand_path = args.get("brand_config")
    if brand_path:
        path = Path(str(brand_path)).expanduser()
    elif DEFAULT_BRAND_PATH.exists():
        path = DEFAULT_BRAND_PATH
    else:
        return None

    if not path.exists():
        return None

    try:
        # Simple TOML-like parsing (avoid dependency)
        brand = {}
        with open(path, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line or line.startswith("#") or line.startswith("["):
                    continue
                if "=" in line:
                    key, _, val = line.partition("=")
                    key = key.strip()
                    val = val.strip().strip('"').strip("'")
                    brand[key] = val
        return brand if brand else None
    except Exception:
        return None


def _apply_slide_master(prs: Presentation, brand: dict, theme: dict) -> None:
    """Apply brand elements to the slide master."""
    logo_path = brand.get("logo", "")
    if not logo_path or not Path(logo_path).expanduser().exists():
        return

    # Add logo to slide master
    try:
        master = prs.slide_masters[0]
        # Add logo to top-right corner
        logo_shape = master.shapes.add_picture(
            str(Path(logo_path).expanduser()),
            Inches(11.5), Inches(0.15), height=Inches(0.45),
        )
    except Exception:
        pass  # Logo insertion is best-effort


def _apply_brand_footer(slide, brand: dict, theme: dict,
                        page_num: int, total: int) -> None:
    """Add brand footer to a slide."""
    footer_text = brand.get("footer", "")
    company = brand.get("company", "")

    y = 7.1
    # Thin accent line
    _add_shape_bg(slide, 0.5, 7.05, 12.3, 0.015, theme["accent"])

    if company:
        _add_textbox(slide, 0.6, y, 4.0, 0.25, company,
                     _make_font(8, theme["muted"]), theme)
    if footer_text:
        _add_textbox(slide, 5.0, y, 6.0, 0.25, footer_text,
                     _make_font(8, theme["muted"]), theme, align=PP_ALIGN.CENTER)
    _add_textbox(slide, 12.0, y, 0.8, 0.25, f"{page_num}/{total}",
                 _make_font(8, theme["muted"]), theme, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
#  AI CONTENT GENERATION
# ═══════════════════════════════════════════════════════════════

AI_SLIDES_PROMPT = """You are a professional presentation designer. Generate a structured slide deck as JSON.

Topic: {topic}
Audience: {audience}
Language: {language}
Style: {style}
Slide count: {slide_count} (content slides, not counting title/TOC/thanks)

Available layouts: title_slide, section_header, content_bullets, two_column,
  image_right, image_left, quote, comparison, timeline, kpi_dashboard,
  team, chart, table, chart_with_text

Available chart types: bar, line, pie, doughnut, area
Available styles: professional, academic, creative, minimal, dark, corporate

Rules:
- First slide should summarize the key message (content_bullets or kpi_dashboard).
- Use kpi_dashboard for numeric metrics (3-4 KPIs with value/label/change).
- Use chart for data trends (provide categories + series with real-looking values).
- Use table for structured comparisons (provide column headers + row data).
- Use comparison for A vs B analysis.
- Use timeline for chronological events (provide date + label per milestone).
- Use section_header to break long decks into parts.
- Use two_column to present balanced perspectives.
- Use quote for testimonials or key insights.
- Every slide must have a title. Bullet points should be concise (5-10 words each).
- Generate realistic, industry-appropriate content (numbers, dates, names).

Return ONLY a JSON object with a "slides" array. Example:
{{
  "slides": [
    {{
      "title": "Market Overview",
      "layout": "kpi_dashboard",
      "kpis": [
        {{"value": "$2.4B", "label": "Market Size 2026", "change": "+18%"}},
        {{"value": "34%", "label": "YoY Growth", "change": "+5%"}}
      ]
    }},
    {{
      "title": "Quarterly Revenue Trend",
      "layout": "chart",
      "chart": {{
        "type": "line",
        "categories": ["Q1", "Q2", "Q3", "Q4"],
        "series": [{{"name": "Revenue", "values": [120, 145, 168, 192]}}]
      }}
    }}
  ]
}}"""


def _ai_generate_slides(args: dict, theme: dict, brand: dict | None) -> list[dict] | None:
    """Use LLM to generate structured slides from topic.

    Requires an OpenAI-compatible client to be available.
    Falls back to template-based generation if LLM is unavailable.
    """
    topic = str(args.get("topic") or args.get("title") or "").strip()
    if not topic:
        return None

    # Try to get a client
    client = _get_llm_client(args)
    if client is None:
        return None

    model = str(args.get("ai_model") or "gpt-4o-mini")
    language = str(args.get("language") or "zh")
    audience = str(args.get("audience") or "")
    style = str(args.get("style") or "professional")
    slide_count = int(args.get("slide_count") or 6)

    prompt = AI_SLIDES_PROMPT.format(
        topic=topic, audience=audience or "general audience",
        language=language, style=style, slide_count=slide_count,
    )

    try:
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "You generate structured presentation content as JSON. Return only valid JSON."},
                {"role": "user", "content": prompt},
            ],
            response_format={"type": "json_object"},
            temperature=0.7,
            max_tokens=4000,
        )
        content = resp.choices[0].message.content
        data = json.loads(content)
        slides = data.get("slides", [])
        if slides and isinstance(slides, list):
            return slides
    except Exception:
        pass

    return None


def _get_llm_client(args: dict):
    """Get an OpenAI-compatible client from args or environment."""
    api_key = args.get("api_key") or os.environ.get("OPENAI_API_KEY")
    base_url = args.get("base_url") or os.environ.get("OPENAI_BASE_URL")

    if not api_key:
        # Try reading from model.toml
        config_path = Path("pysrc/model.toml")
        if config_path.exists():
            try:
                import tomllib
            except ImportError:
                import tomli as tomllib
            try:
                with open(config_path, "rb") as f:
                    cfg = tomllib.load(f)
                llm = cfg.get("llm", {})
                active = llm.get("active_provider", "openai")
                provider = llm.get(active, {})
                api_key = api_key or provider.get("api_key", "")
                base_url = base_url or provider.get("base_url", "")
            except Exception:
                pass

    if not api_key:
        return None

    try:
        from openai import OpenAI
        kwargs = {"api_key": api_key}
        if base_url:
            kwargs["base_url"] = base_url
        return OpenAI(**kwargs)
    except ImportError:
        return None


# ═══════════════════════════════════════════════════════════════
#  EXTERNAL ENGINE — full pipeline integration
# ═══════════════════════════════════════════════════════════════


def _resolve_external_skill_dir(args: dict) -> Path | None:
    candidates = [
        args.get("ppt_master_path"),
        os.environ.get("PPT_MASTER_PATH"),
        os.environ.get("PPT_MASTER_SKILL_DIR"),
        str(DEFAULT_EXTERNAL_ROOT),
        str(Path(".runtime") / "ppt-master"),
    ]
    for raw in candidates:
        if not raw:
            continue
        base = Path(str(raw)).expanduser()
        possible = [base, base / "skills" / "ppt-master"]
        for item in possible:
            if (item / "scripts").is_dir() and (item / "SKILL.md").exists():
                return item.resolve()
    return None


def _target_external_root(args: dict) -> Path:
    raw = args.get("ppt_master_path") or os.environ.get("PPT_MASTER_PATH") or DEFAULT_EXTERNAL_ROOT
    path = Path(str(raw)).expanduser()
    if path.name == "ppt-master" and path.parent.name == "skills":
        return path.parent.parent
    return path


def _external_repo_root(skill_dir: Path) -> Path:
    try:
        if skill_dir.parts[-2:] == ("skills", "ppt-master"):
            return skill_dir.parent.parent
    except Exception:
        pass
    return skill_dir


def _script_path(skill_dir: Path, relative: str) -> Path:
    return skill_dir / Path(relative)


def check_external_engine(args: dict) -> dict:
    skill_dir = _resolve_external_skill_dir(args)
    timeout = int(args.get("timeout_seconds") or 300)
    python_exe = _python_executable(args)
    if not skill_dir:
        return {
            "status": "error", "installed": False,
            "message": "External PPT Master is not installed. Run action=setup.",
            "default_path": str(DEFAULT_EXTERNAL_ROOT),
        }
    repo_root = _external_repo_root(skill_dir)
    scripts = [
        {"path": str((skill_dir / rel).resolve()), "exists": (skill_dir / rel).exists()}
        for rel in REQUIRED_EXTERNAL_SCRIPTS
    ]
    # Check for full pipeline availability
    pipeline_scripts = [
        "scripts/generate_slides.py", "scripts/render_slides.py",
        "scripts/apply_template.py", "scripts/export_pptx.py",
    ]
    pipeline = [
        {"path": str((skill_dir / rel).resolve()), "exists": (skill_dir / rel).exists()}
        for rel in pipeline_scripts
    ]
    return {
        "status": "success", "installed": True,
        "ppt_master_path": str(skill_dir),
        "repo_root": str(repo_root),
        "scripts": scripts,
        "pipeline": pipeline,
        "ready": all(item["exists"] for item in scripts),
        "pipeline_ready": all(item["exists"] for item in pipeline),
    }


def setup_external_engine(args: dict) -> dict:
    timeout = int(args.get("timeout_seconds") or 300)
    python_exe = _python_executable(args)
    target_root = _target_external_root(args).resolve()
    logs: list[dict] = []

    existing = _resolve_external_skill_dir({"ppt_master_path": str(target_root)})
    if existing:
        repo_root = _external_repo_root(existing)
        if _truthy(args.get("update"), False) and (repo_root / ".git").exists():
            proc = _run_command(["git", "pull", "--ff-only"], cwd=repo_root, timeout_seconds=timeout)
            logs.append({"step": "git_pull", "ok": proc.returncode == 0})
            if proc.returncode != 0:
                return {"status": "error", "message": "git pull failed.", "logs": logs}
    else:
        if target_root.exists() and any(target_root.iterdir()):
            return {"status": "error", "message": f"Target path exists but is not PPT Master: {target_root}"}
        target_root.parent.mkdir(parents=True, exist_ok=True)
        proc = _run_command(["git", "clone", "--depth", "1", UPSTREAM_REPO_URL, str(target_root)],
                           cwd=None, timeout_seconds=timeout)
        logs.append({"step": "git_clone", "ok": proc.returncode == 0})
        if proc.returncode != 0:
            return {"status": "error", "message": "git clone failed.", "logs": logs}

    skill_dir = _resolve_external_skill_dir({"ppt_master_path": str(target_root)})
    if not skill_dir:
        return {"status": "error", "message": "Clone completed but skills/ppt-master not found."}

    repo_root = _external_repo_root(skill_dir)
    requirements = next((p for p in [repo_root / "requirements.txt", skill_dir / "requirements.txt"] if p.exists()), None)
    if _truthy(args.get("install_dependencies"), True) and requirements:
        proc = _run_command([python_exe, "-m", "pip", "install", "-r", str(requirements)],
                           cwd=repo_root, timeout_seconds=timeout)
        logs.append({"step": "pip_install", "ok": proc.returncode == 0})
        if proc.returncode != 0:
            return {"status": "error", "message": "Dependency installation failed."}

    health = check_external_engine({**args, "ppt_master_path": str(skill_dir)})
    return {"status": "success" if health.get("ready") else "error",
            "message": "External PPT Master is ready." if health.get("ready") else "Health check incomplete.",
            "ppt_master_path": str(skill_dir), "repo_root": str(repo_root), "logs": logs, "health": health}


def _run_external_pipeline(args: dict, skill_dir: Path, output_path: Path) -> dict:
    """Run the full external PPT Master pipeline: generate → render → export.

    This uses the external engine's complete pipeline (not just file conversion):
      1. scripts/generate_slides.py — AI content generation from topic/outline
      2. scripts/render_slides.py — SVG rendering with templates
      3. scripts/finalize_svg.py — SVG post-processing
      4. scripts/svg_to_pptx.py — Native PPTX export
    """
    timeout = int(args.get("timeout_seconds") or 600)
    python_exe = _python_executable(args)
    work_dir = (Path(args.get("work_dir") or "outputs/ppt_master/external_work")).expanduser().resolve()
    work_dir.mkdir(parents=True, exist_ok=True)
    results: list[dict] = []

    # Step 1: Generate slides content via AI (if topic/outline provided)
    generate_script = _script_path(skill_dir, "scripts/generate_slides.py")
    slides_json = work_dir / "slides.json"
    if generate_script.exists():
        topic = str(args.get("topic") or args.get("title") or "Presentation")
        outline = str(args.get("outline") or "")
        command = [python_exe, str(generate_script), "--topic", topic,
                   "--output", str(slides_json),
                   "--style", str(args.get("style") or "professional"),
                   "--language", str(args.get("language") or "zh")]
        if outline:
            command.extend(["--outline", outline])
        proc = _run_command(command, cwd=skill_dir, timeout_seconds=timeout)
        results.append({"step": "generate", "ok": proc.returncode == 0, "stdout": proc.stdout[-2000:]})
        if proc.returncode == 0 and slides_json.exists():
            args = dict(args)
            try:
                args["slides"] = json.loads(slides_json.read_text(encoding="utf-8"))
            except Exception:
                pass
        else:
            # generate_slides.py not available or failed — use builtin content gen
            results.append({"step": "generate", "ok": False,
                           "message": "External generate not available, using builtin content"})

    # Step 2: Render slides to SVG
    render_script = _script_path(skill_dir, "scripts/render_slides.py")
    svg_dir = work_dir / "svg"
    svg_dir.mkdir(exist_ok=True)
    if render_script.exists():
        command = [python_exe, str(render_script),
                   "--slides", str(slides_json) if slides_json.exists() else json.dumps(_normalize_slides(args)),
                   "--output", str(svg_dir),
                   "--style", str(args.get("style") or "professional"),
                   "--theme-color", str(args.get("theme_color") or theme_from_style(args).get("accent", "#2563eb"))]
        proc = _run_command(command, cwd=skill_dir, timeout_seconds=timeout)
        results.append({"step": "render", "ok": proc.returncode == 0})
    else:
        results.append({"step": "render", "ok": False, "message": "render_slides.py not found"})

    # Step 3: Finalize SVGs
    finalize_script = _script_path(skill_dir, "scripts/finalize_svg.py")
    if finalize_script.exists():
        proc = _run_command([python_exe, str(finalize_script), str(svg_dir)],
                           cwd=skill_dir, timeout_seconds=timeout)
        results.append({"step": "finalize", "ok": proc.returncode == 0})

    # Step 4: Export to PPTX
    export_script = _script_path(skill_dir, "scripts/svg_to_pptx.py")
    if export_script.exists() and any(svg_dir.glob("*.svg")):
        _ensure_output_writable(output_path, _truthy(args.get("overwrite"), False))
        command = [python_exe, str(export_script), str(svg_dir),
                   "-o", str(output_path), "--only", "native"]
        if args.get("template"):
            command.extend(["--template", str(args["template"])])
        proc = _run_command(command, cwd=skill_dir, timeout_seconds=timeout)
        results.append({"step": "export", "ok": proc.returncode == 0 and output_path.exists()})
        if output_path.exists():
            return {"status": "success", "format": "pptx", "output": str(output_path.resolve()),
                    "editable": True, "engine": "external_full_pipeline",
                    "ppt_master_path": str(skill_dir), "pipeline_results": results}

    # Fallback: use builtin renderer
    builtin_result = build_presentation(args, engine="builtin_fallback")
    builtin_result["external_pipeline_results"] = results
    builtin_result["external_pipeline_status"] = "partial"
    return builtin_result


def theme_from_style(args: dict) -> dict:
    style = str(args.get("style") or "professional").strip().lower()
    return THEMES.get(style, THEMES["professional"])


def list_templates(args: dict) -> dict:
    skill_dir = _resolve_external_skill_dir(args)
    external = []
    if skill_dir:
        for path in [
            skill_dir / "templates" / "layouts" / "layouts_index.json",
            skill_dir / "templates" / "styles" / "styles_index.json",
            skill_dir / "templates" / "themes" / "themes_index.json",
        ]:
            if path.exists():
                try:
                    payload = json.loads(path.read_text(encoding="utf-8"))
                    if isinstance(payload, list):
                        external.extend(payload)
                    elif isinstance(payload, dict):
                        for key in ("layouts", "styles", "themes", "templates"):
                            items = payload.get(key)
                            if isinstance(items, list):
                                external.extend(items)
                except Exception:
                    continue
    return {
        "status": "success", "action": "list_templates",
        "engine": "external" if external else "builtin",
        "templates": external or TEMPLATE_CATALOG,
    }


# ═══════════════════════════════════════════════════════════════
#  ENTRY POINT
# ═══════════════════════════════════════════════════════════════

def main(**kwargs) -> dict:
    args = dict(kwargs) if kwargs else parse_cli_args()
    try:
        action = str(args.get("action") or "generate").strip().lower()
        if action in {"list_templates", "templates"}:
            return list_templates(args)
        if action in {"setup", "install"}:
            return setup_external_engine(args)
        if action in {"check", "health", "doctor"}:
            return check_external_engine(args)
        if action not in {"generate", "create", "render"}:
            return {"status": "error", "error": f"Unknown action: {action}"}

        title = str(args.get("title") or args.get("topic") or "Presentation").strip()
        mode = str(args.get("mode") or "auto").strip().lower()
        if mode not in {"auto", "external", "builtin"}:
            return {"status": "error", "error": f"Unknown mode: {mode}"}

        has_external_input = bool(
            args.get("source") or args.get("project") or
            args.get("external_project") or args.get("svg_dir")
        )

        # mode=builtin or auto with no external input → builtin
        if mode == "builtin" or (mode == "auto" and not has_external_input):
            return build_presentation(args, engine="builtin")

        # mode=external or auto with external available → try full pipeline
        skill_dir = _resolve_external_skill_dir(args)
        if not skill_dir and _truthy(args.get("auto_install"), True):
            setup = setup_external_engine(args)
            if setup.get("status") == "success":
                skill_dir = _resolve_external_skill_dir({"ppt_master_path": setup.get("ppt_master_path")})

        output_path = _default_output_path(args, title)

        if skill_dir:
            # Check if full pipeline is available
            health = check_external_engine({**args, "ppt_master_path": str(skill_dir)})
            if health.get("pipeline_ready"):
                result = _run_external_pipeline(args, skill_dir, output_path)
                if result.get("status") == "success":
                    return result

            # Fall back to source conversion + builtin render
            source = str(args.get("source") or "").strip()
            if source:
                conversion = _convert_source_to_markdown(args, skill_dir, source)
                if conversion.get("status") == "success":
                    args = dict(args)
                    args["outline"] = conversion["markdown"]
                    result = build_presentation(args, engine="external_source_conversion_builtin_render")
                    result["ppt_master_path"] = str(skill_dir)
                    result["external_markdown_path"] = conversion.get("markdown_path")
                    return result

        if mode == "external":
            requested_path = args.get("ppt_master_path") or os.environ.get("PPT_MASTER_PATH") or DEFAULT_EXTERNAL_ROOT
            return {
                "status": "error",
                "message": f"External PPT Master path is not available or pipeline failed: {requested_path}",
            }

        # Fallback
        return build_presentation(args, engine="builtin_fallback")

    except Exception as exc:
        return {"status": "error", "error": str(exc)}


# ═══════════════════════════════════════════════════════════════
#  SOURCE CONVERSION (kept from original for backward compat)
# ═══════════════════════════════════════════════════════════════

def _convert_source_to_markdown(args: dict, skill_dir: Path, source: str) -> dict:
    candidates = [
        _script_path(skill_dir, rel)
        for suffix, scripts in SOURCE_CONVERTERS.items()
        if source.lower().endswith(suffix) or _is_url(source)
        for rel in scripts
    ]
    candidates = [p for p in candidates if p.exists()]
    if not candidates:
        return {"status": "error", "message": f"No external source converter found for {source}."}

    out_dir = (Path(args.get("work_dir") or "outputs/ppt_master/external_work")).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    out_md = out_dir / f"{_safe_filename(Path(source).stem if not _is_url(source) else 'web_source')}.md"
    timeout = int(args.get("timeout_seconds") or 300)

    for script in candidates:
        command = [_python_executable(args), str(script), str(source), "-o", str(out_md)]
        proc = _run_command(command, cwd=skill_dir, timeout_seconds=timeout)
        if proc.returncode == 0 and out_md.exists():
            return {"status": "success",
                    "markdown": out_md.read_text(encoding="utf-8", errors="replace"),
                    "markdown_path": str(out_md.resolve())}
    return {"status": "error", "message": "External source conversion failed."}


if __name__ == "__main__":
    result = main()
    print(json.dumps(result, ensure_ascii=False, indent=2))
    raise SystemExit(0 if result.get("status") == "success" else 2)
