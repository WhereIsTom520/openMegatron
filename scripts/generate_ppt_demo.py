"""
Generate demo PPT showcasing all enhanced layouts.
"""
import sys, importlib.util, os

spec = importlib.util.spec_from_file_location(
    'enhanced_layouts',
    'pysrc/skills/office/ppt-master-1.0.0/scripts/enhanced_layouts.py'
)
enhanced = importlib.util.module_from_spec(spec)
spec.loader.exec_module(enhanced)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

theme = {
    'accent': '#6366f1', 'dark': '#1e293b', 'muted': '#64748b',
    'light': '#f1f5f9', 'background': '#ffffff',
    'secondary': '#14b8a6', 'success': '#10b981', 'warning': '#f59e0b',
    'font_title': 'Calibri', 'font_body': 'Calibri',
}

def _hex_rgb(h):
    h = h.lstrip('#')
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _add_bg(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(*_hex_rgb(c))
    s.line.fill.background()
    return s

def _txt(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    f = box.text_frame
    f.word_wrap = True
    p = f.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.font.size = Pt(size)
    r.font.color.rgb = RGBColor(*_hex_rgb(color))
    r.font.bold = bold
    r.font.name = 'Calibri'
    return box

# === Slide 1: Title ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, 0, 0, 13.333, 7.5, theme['background'])
_add_bg(slide, 0, 0, 0.5, 7.5, theme['accent'])
_txt(slide, 1.2, 1.8, 11.0, 1.5, 'OpenMegatron PPT Master', 48, theme['dark'], True)
_txt(slide, 1.25, 3.3, 10.5, 0.8, 'Enhanced Layout System - Build Professional Decks in Seconds', 20, theme['muted'])
_add_bg(slide, 1.25, 4.3, 2.5, 0.05, theme['accent'])
_txt(slide, 1.25, 5.0, 5.0, 0.4, 'AI-Powered Presentation Engine', 14, theme['dark'])
_txt(slide, 1.25, 5.4, 5.0, 0.4, '21 Layout Types | 6 Theme Styles | Native PPTX Output', 12, theme['muted'])
_txt(slide, 1.25, 6.5, 5.0, 0.3, 'June 2026', 11, theme['muted'])

# === Slide 2: Code Block ===
enhanced.build_code_block(prs, {
    'title': 'Code Block - Syntax-Aware Presentation',
    'language': 'python',
    'code': (
        'import torch\n'
        'import torch.nn as nn\n'
        '\n'
        'class TransformerBlock(nn.Module):\n'
        '    """Multi-head self-attention with FFN."""\n'
        '    def __init__(self, dim=512, heads=8, dropout=0.1):\n'
        '        super().__init__()\n'
        '        self.attention = nn.MultiheadAttention(dim, heads)\n'
        '        self.ffn = nn.Sequential(\n'
        '            nn.Linear(dim, dim * 4),\n'
        '            nn.GELU(),\n'
        '            nn.Linear(dim * 4, dim),\n'
        '        )\n'
        '        self.norm1 = nn.LayerNorm(dim)\n'
        '        self.norm2 = nn.LayerNorm(dim)\n'
        '    \n'
        '    def forward(self, x):\n'
        '        # Self-attention with residual\n'
        '        attn_out, _ = self.attention(x, x, x)\n'
        '        x = self.norm1(x + attn_out)\n'
        '        # Feed-forward with residual\n'
        '        return self.norm2(x + self.ffn(x))'
    ),
    'caption': 'PyTorch Transformer implementation - production-ready code'
}, 2, theme)

# === Slide 3: Flowchart Vertical ===
enhanced.build_flowchart(prs, {
    'title': 'AI Agent Pipeline - From Query to Response',
    'steps': [
        {'title': 'Intent Parsing', 'description': 'Analyze user query, extract entities and intent using NLP'},
        {'title': 'Task Decomposition', 'description': 'Break complex request into sub-tasks with dependency graph'},
        {'title': 'Skill Routing', 'description': 'Match sub-tasks to appropriate skills and model tiers'},
        {'title': 'Execution & Validation', 'description': 'Run tools, validate outputs, handle errors with retry'},
        {'title': 'Response Synthesis', 'description': 'Merge results, format citations, ensure factual accuracy'},
    ]
}, 3, theme)

# === Slide 4: Three Column ===
enhanced.build_three_column(prs, {
    'title': 'LLM Provider Comparison - Choose the Right Model',
    'title1': 'GPT-4o',
    'title2': 'External Agent advanced cloud model 4.8',
    'title3': 'Gemini 2.5 Pro',
    'column1': [
        'Latest GPT generation',
        '128K context window',
        'Native JSON mode',
        'Function calling',
        'Vision + audio input',
        'Affordable pricing',
    ],
    'column2': [
        'Deep chain-of-thought',
        '200K context window',
        'Tool use (MCP native)',
        'Code generation & analysis',
        'Long-form document understanding',
        'Enterprise-grade safety',
    ],
    'column3': [
        'Google Cloud ecosystem',
        '1M+ context window',
        'Multimodal by design',
        'Vertex AI integration',
        'Real-time web grounding',
        'Competitive free tier',
    ],
}, 4, theme)

# === Slide 5: Grid 2x2 ===
enhanced.build_grid_2x2(prs, {
    'title': 'Core Capabilities - What Makes It Powerful',
    'cards': [
        {'title': 'Multi-Model Dispatch', 'description': 'Auto-route tasks to optimal models (Lite/Standard/Advanced) based on capability probing and cost analysis.', 'icon': '1'},
        {'title': 'Hybrid RAG Pipeline', 'description': 'Tri-store architecture combining PostgreSQL vector search, Neo4j graph traversal, and Redis semantic caching.', 'icon': '2'},
        {'title': 'Self-Evolution Engine', 'description': 'Continuous learning from task trajectories. Companion model replaces cloud API for simple tasks over time.', 'icon': '3'},
        {'title': 'Skill Marketplace', 'description': '30+ pluggable skills across code, research, office, media, and agent domains. Add custom skills via SKILL.md.', 'icon': '4'},
    ]
}, 5, theme, True)

# === Slide 6: Comparison Plus ===
enhanced.build_comparison_plus(prs, {
    'title': 'OpenMegatron vs Traditional AI Platforms',
    'left_title': 'Traditional AI Platforms',
    'right_title': 'OpenMegatron',
    'features': [
        {'name': 'Self-hosted / On-premise', 'left': 'x', 'right': 'yes'},
        {'name': 'Multi-model Routing', 'left': 'x', 'right': 'yes'},
        {'name': 'Local Companion Model', 'left': 'x', 'right': 'yes'},
        {'name': 'Research Workflow Engine', 'left': 'x', 'right': 'yes'},
        {'name': 'Anti-hallucination Layer', 'left': 'x', 'right': 'yes'},
        {'name': 'PRISMA Systematic Review', 'left': 'x', 'right': 'yes'},
        {'name': 'Open Source', 'left': 'Partial', 'right': 'yes'},
    ]
}, 6, theme)

# === Slide 7: Flowchart Horizontal ===
enhanced.build_flowchart(prs, {
    'title': 'Development Workflow - From Idea to Production',
    'direction': 'horizontal',
    'steps': [
        {'title': 'Ideate', 'description': 'Define problem'},
        {'title': 'Prototype', 'description': 'Build MVP with AI'},
        {'title': 'Validate', 'description': 'Tests & benchmarks'},
        {'title': 'Deploy', 'description': 'Ship to staging'},
        {'title': 'Monitor', 'description': 'Track & iterate'},
    ]
}, 7, theme)

# === Slide 8: KPI Dashboard ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, 0.6, 0.35, 0.06, 0.55, theme['accent'])
_txt(slide, 0.78, 0.32, 11.5, 0.7, 'Performance Metrics - Impact at a Glance', 28, theme['dark'], True)
kpis = [
    {'value': '21', 'label': 'Layout Types', 'change': '+5 (New)'},
    {'value': '6', 'label': 'Theme Styles', 'change': '+1'},
    {'value': '100%', 'label': 'Native Editable', 'change': 'Industry Best'},
    {'value': '<3s', 'label': 'Avg Generation', 'change': 'Ultra Fast'},
]
n = len(kpis)
card_w = 11.5 / n - 0.3
for i, kpi in enumerate(kpis):
    x = 0.8 + i * (card_w + 0.3)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(card_w), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*_hex_rgb(theme['light']))
    card.line.fill.background()
    _txt(slide, x + 0.2, 1.8, card_w - 0.4, 1.0, kpi['value'], 40, theme['accent'], True, PP_ALIGN.CENTER)
    _txt(slide, x + 0.2, 2.9, card_w - 0.4, 0.5, kpi['label'], 14, theme['dark'], align=PP_ALIGN.CENTER)
    _txt(slide, x + 0.2, 3.4, card_w - 0.4, 0.4, kpi['change'], 12, theme['success'], True, PP_ALIGN.CENTER)

# === Slide 9: Thank You ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, 0, 0, 13.333, 7.5, theme['accent'])
_txt(slide, 1.5, 2.2, 10.3, 1.5, 'Thank You', 56, '#ffffff', True, PP_ALIGN.CENTER)
_txt(slide, 1.5, 3.8, 10.3, 0.6, 'OpenMegatron PPT Master - Enhanced Layout System', 18, '#ffffffcc', align=PP_ALIGN.CENTER)
_add_bg(slide, 5.5, 4.7, 2.3, 0.04, '#ffffff88')
_txt(slide, 1.5, 5.2, 10.3, 0.8, 'github.com/WhereIsTom520/openMegatron', 14, '#ffffff99', align=PP_ALIGN.CENTER)

# === SAVE ===
output_path = 'outputs/openmegatron_ppt_master_demo.pptx'
os.makedirs('outputs', exist_ok=True)
prs.save(output_path)
file_size = os.path.getsize(output_path)

print(f'Generated: {output_path}')
print(f'Size: {file_size} bytes ({file_size/1024:.1f} KB)')
print(f'Slides: {len(prs.slides)}')
layouts = ['Title', 'Code Block', 'Flowchart (V)', 'Three Column', 'Grid 2x2', 'Comparison Plus', 'Flowchart (H)', 'KPI Dashboard', 'Thank You']
for i, name in enumerate(layouts):
    print(f'  Slide {i+1}: {name} ({len(prs.slides[i].shapes)} shapes)')
